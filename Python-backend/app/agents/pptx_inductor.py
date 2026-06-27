"""
PPTX 模板归纳流水线 — 从参考 PPTX 中自动学习布局结构

流程:
  1. 功能页识别（LLM 分类 opening/toc/section/ending/content）
  2. 结构指纹提取 + 聚类（替代 PPTAgent 的 ViT 嵌入）
  3. LLM Schema 提取（为每个聚类命名元素）
  4. 输出 slide_induction.json
"""
import json
import logging
import math
from collections import defaultdict
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger("agents.pptx_inductor")


# ==================== 1. 幻灯片文本提取 ====================

def _extract_slide_text(slide) -> Tuple[str, str]:
    """提取幻灯片的标题和正文文本"""
    title = ""
    body_parts = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = "\n".join(p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip())
        if not text:
            continue
        # 简单启发式：最高/最大的文字形状当作标题
        if not title and shape.top and shape.top < Emu(2000000):
            title = text[:100]
        else:
            body_parts.append(text[:200])
    return title, "\n".join(body_parts)[:300]


def _extract_all_slides_text(prs) -> List[dict]:
    """提取所有幻灯片的文本摘要"""
    slides_text = []
    for i, slide in enumerate(prs.slides):
        title, body = _extract_slide_text(slide)
        slides_text.append({
            "index": i,
            "slide_number": i + 1,
            "title": title,
            "body": body,
        })
    return slides_text


# ==================== 2. 功能页识别（LLM） ====================

def _classify_functional_slides(slides_text: List[dict], llm) -> Tuple[dict, List[int]]:
    """用 LLM 批量分类幻灯片为功能页或内容页

    Returns:
        functional_map: {"opening": [1], "table of contents": [2], "ending": [12], ...}
        content_indices: 内容页的 0-based 索引列表
    """
    slides_desc = []
    for s in slides_text:
        title = s["title"] or "(无标题)"
        body_preview = s["body"][:80] if s["body"] else ""
        slides_desc.append(f"Slide {s['slide_number']}: [{title}] {body_preview}")

    prompt = f"""你是一个 PPT 结构分析专家。请判断以下每张幻灯片的功能类型。

幻灯片列表:
{chr(10).join(slides_desc)}

请返回 JSON 对象，键为功能类型，值为幻灯片编号列表（1-based）。
功能类型包括:
- "opening": 封面/开场页（通常包含演讲标题、作者名、日期）
- "table of contents": 目录页（列出各章节/部分）
- "section outline": 章节分隔页（标识新章节的开始，通常只有章节标题）
- "ending": 结尾页（致谢、Q&A、联系方式等）
- "content": 正文内容页（包含具体知识点、论述、图表等）

注意: 一张幻灯片只能归入一个类型。如果不确定，归为 "content"。

只返回 JSON，不要其他文字。格式示例:
{{"opening": [1], "table of contents": [2], "section outline": [5, 9], "content": [3, 4, 6, 7, 8, 10, 11], "ending": [12]}}"""

    try:
        response = llm.invoke(prompt)
        response_text = response.content if hasattr(response, "content") else str(response)
        # 提取 JSON
        import re
        json_match = re.search(r'\{[^{}]*\}', response_text, re.DOTALL)
        if json_match:
            result = json.loads(json_match.group())
        else:
            result = json.loads(response_text)
    except Exception as e:
        logger.warning(f"[归纳] 功能页分类 LLM 调用失败，全部视为内容页: {e}")
        result = {}

    functional_map = {}
    content_indices = []
    all_classified = set()

    for func_type, slide_numbers in result.items():
        if func_type == "content":
            continue
        functional_map[func_type] = slide_numbers
        all_classified.update(slide_numbers)

    for s in slides_text:
        if s["slide_number"] not in all_classified:
            content_indices.append(s["index"])

    return functional_map, content_indices


# ==================== 3. 结构指纹提取 ====================

def _calc_text_area_ratio(slide) -> float:
    """计算文本形状面积占幻灯片总面积的比例"""
    slide_area = slide.part.slide_layout.slide_master.part.presentation.slide_width * \
                 slide.part.slide_layout.slide_master.part.presentation.slide_height
    if slide_area == 0:
        return 0.0
    text_area = 0
    for shape in slide.shapes:
        if shape.has_text_frame and shape.width and shape.height:
            text_area += shape.width * shape.height
    return min(text_area / slide_area, 1.0)


def _count_bullets(slide) -> int:
    """统计幻灯片中的项目符号数量"""
    count = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text and (text.startswith("•") or text.startswith("-") or text.startswith("·") or
                         (len(text) > 2 and text[0].isdigit() and text[1] in ".、)")):
                count += 1
    return count


def _get_max_font_size(slide) -> float:
    """获取幻灯片中最大的字体大小（pt）"""
    max_size = 0.0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    size_pt = run.font.size / 12700  # EMU to pt
                    if size_pt > max_size:
                        max_size = size_pt
    return max_size


def _structural_fingerprint(slide) -> dict:
    """提取幻灯片结构指纹（用于聚类）"""
    shape_count = 0
    text_shape_count = 0
    rect_count = 0
    picture_count = 0
    table_count = 0
    chart_count = 0

    for shape in slide.shapes:
        shape_count += 1
        if shape.has_text_frame:
            text_shape_count += 1
        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.RECTANGLE or st == MSO_SHAPE_TYPE.ROUNDED_RECTANGLE:
            rect_count += 1
        elif st == MSO_SHAPE_TYPE.PICTURE:
            picture_count += 1
        elif st == MSO_SHAPE_TYPE.TABLE:
            table_count += 1
        elif st == MSO_SHAPE_TYPE.CHART:
            chart_count += 1

    return {
        "shape_count": shape_count,
        "text_shape_count": text_shape_count,
        "rect_count": rect_count,
        "picture_count": picture_count,
        "table_count": table_count,
        "chart_count": chart_count,
        "text_area_ratio": round(_calc_text_area_ratio(slide), 3),
        "bullet_count": _count_bullets(slide),
        "max_font_size": round(_get_max_font_size(slide), 1),
    }


def _fingerprint_distance(fp1: dict, fp2: dict) -> float:
    """计算两个结构指纹之间的归一化欧氏距离"""
    keys = ["shape_count", "text_shape_count", "rect_count", "picture_count",
            "table_count", "chart_count", "text_area_ratio", "bullet_count", "max_font_size"]
    # 各维度的典型范围（用于归一化）
    ranges = {
        "shape_count": 20, "text_shape_count": 10, "rect_count": 10,
        "picture_count": 5, "table_count": 3, "chart_count": 3,
        "text_area_ratio": 1.0, "bullet_count": 15, "max_font_size": 40,
    }
    total = 0.0
    for k in keys:
        r = ranges.get(k, 1.0)
        if r == 0:
            continue
        diff = (fp1.get(k, 0) - fp2.get(k, 0)) / r
        total += diff * diff
    return math.sqrt(total)


# ==================== 4. 结构聚类 ====================

def _cluster_content_slides(prs, content_indices: List[int], threshold: float = 0.35) -> List[List[int]]:
    """对内容页按结构指纹进行凝聚聚类

    Returns:
        clusters: 每个聚类包含的 0-based slide 索引列表
    """
    if not content_indices:
        return []

    # 提取指纹
    fingerprints = {}
    for idx in content_indices:
        slide = prs.slides[idx]
        fingerprints[idx] = _structural_fingerprint(slide)

    # 先按 slide_layout.name 分组（如果有的话）
    layout_groups = defaultdict(list)
    for idx in content_indices:
        slide = prs.slides[idx]
        layout_name = slide.slide_layout.name if slide.slide_layout else "default"
        layout_groups[layout_name].append(idx)

    clusters = []
    for layout_name, indices in layout_groups.items():
        if len(indices) <= 1:
            clusters.append(indices)
            continue

        # 凝聚聚类
        sub_clusters = [[idx] for idx in indices]

        while len(sub_clusters) > 1:
            min_dist = float("inf")
            merge_i, merge_j = -1, -1

            for i in range(len(sub_clusters)):
                for j in range(i + 1, len(sub_clusters)):
                    # 用两个聚类中代表元素的距离（取平均）
                    dists = []
                    for a in sub_clusters[i]:
                        for b in sub_clusters[j]:
                            dists.append(_fingerprint_distance(fingerprints[a], fingerprints[b]))
                    avg_dist = sum(dists) / len(dists)
                    if avg_dist < min_dist:
                        min_dist = avg_dist
                        merge_i, merge_j = i, j

            if min_dist > threshold:
                break

            sub_clusters[merge_i] = sub_clusters[merge_i] + sub_clusters[merge_j]
            sub_clusters.pop(merge_j)

        clusters.extend(sub_clusters)

    return clusters


# ==================== 5. Schema 提取（LLM） ====================

def _get_shape_preview(shape, max_len: int = 60) -> str:
    """获取形状的文本预览"""
    if not shape.has_text_frame:
        return ""
    texts = []
    for p in shape.text_frame.paragraphs:
        t = p.text.strip()
        if t:
            texts.append(t)
    full = " | ".join(texts)
    return full[:max_len] + ("..." if len(full) > max_len else "")


def _extract_schema(cluster_indices: List[int], prs, llm) -> dict:
    """从一个聚类的模板幻灯片中提取元素 schema

    Returns:
        {"layout_name:subtype": {"template_id": int, "slides": [...], "elements": [...]}}
    """
    # 用第一张幻灯片作为代表
    representative_idx = cluster_indices[0]
    slide = prs.slides[representative_idx]

    shapes_info = []
    for i, shape in enumerate(slide.shapes):
        pos = f"(left={shape.left}, top={shape.top}, w={shape.width}, h={shape.height})"
        preview = _get_shape_preview(shape)
        shape_type = "text" if shape.has_text_frame else (
            "image" if shape.shape_type == MSO_SHAPE_TYPE.PICTURE else
            "table" if shape.shape_type == MSO_SHAPE_TYPE.TABLE else
            "chart" if shape.shape_type == MSO_SHAPE_TYPE.CHART else "other"
        )
        shapes_info.append(f"Shape {i} [{shape_type}] {pos}: {preview}")

    prompt = f"""你是一个 PPT 模板分析专家。请分析以下幻灯片的形状布局，为每个有意义的内容元素命名。

幻灯片形状信息:
{chr(10).join(shapes_info)}

请返回一个 JSON 对象，包含:
1. "layout_name": 布局的简短描述性英文名称（如 "Title with Bulleted List"、"Two Column Comparison"、"Image with Caption"），后面加冒号和子类型（:text 或 :image）
2. "elements": 元素列表，每个元素包含:
   - "name": 元素的语义名称（小写英文，如 "main title"、"main bullets"、"main image"、"quote and explanation"）
   - "shape_idx": 对应的 Shape 索引号
   - "type": "text" 或 "image"

规则:
- 只命名有意义的内容元素，忽略装饰性形状（如背景矩形、分隔线）
- 标题类元素命名为 "main title" 或 "section title"
- 主要文本块命名为 "main bullets"、"main paragraph"、"quote and explanation" 等
- 图片元素命名为 "main image"
- 如果布局有多个文本区域，用 "left section"、"right section" 区分
- 子类型: 如果主要是文字布局用 :text，如果包含重要图片用 :image

只返回 JSON，格式示例:
{{"layout_name": "Title with Bulleted List:text", "elements": [{{"name": "main title", "shape_idx": 0, "type": "text"}}, {{"name": "main bullets", "shape_idx": 1, "type": "text"}}]}}"""

    try:
        response = llm.invoke(prompt)
        response_text = response.content if hasattr(response, "content") else str(response)
        import re
        json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
        if json_match:
            schema = json.loads(json_match.group())
        else:
            schema = json.loads(response_text)
    except Exception as e:
        logger.warning(f"[归纳] Schema 提取失败 (slide {representative_idx + 1}): {e}")
        # 降级：用默认名称
        schema = {
            "layout_name": f"Content Layout {representative_idx + 1}:text",
            "elements": [],
        }
        for i, shape in enumerate(slide.shapes):
            if shape.has_text_frame and _get_shape_preview(shape):
                schema["elements"].append({
                    "name": f"text_{i}",
                    "shape_idx": i,
                    "type": "text",
                })

    # 构建 induction 格式
    layout_name = schema.get("layout_name", f"Unknown Layout:text")
    elements = []
    for elem in schema.get("elements", []):
        shape_idx = elem.get("shape_idx", 0)
        if shape_idx < len(list(slide.shapes)):
            shape = list(slide.shapes)[shape_idx]
            preview = _get_shape_preview(shape, max_len=200)
            elements.append({
                "name": elem["name"],
                "type": elem.get("type", "text"),
                "data": [preview] if preview else [],
            })

    return {
        layout_name: {
            "template_id": representative_idx + 1,  # 1-based
            "slides": [idx + 1 for idx in cluster_indices],  # 1-based
            "elements": elements,
        }
    }


# ==================== 6. 主流水线 ====================

def induct_template(pptx_path: str, llm=None) -> dict:
    """从 PPTX 模板中归纳布局结构，生成 slide_induction.json

    Args:
        pptx_path: PPTX 文件路径
        llm: LLM 实例（可选，不传则自动创建）

    Returns:
        与现有 slide_induction.json 格式一致的 dict
    """
    from app.agents.llm_factory import MIMOClient

    if llm is None:
        llm = MIMOClient(temperature=0.3, max_tokens=4096)

    logger.info(f"[归纳] 开始归纳模板: {pptx_path}")
    prs = Presentation(pptx_path)
    total_slides = len(prs.slides)
    logger.info(f"[归纳] 共 {total_slides} 张幻灯片")

    # 1. 提取文本
    slides_text = _extract_all_slides_text(prs)

    # 2. 功能页识别
    functional_map, content_indices = _classify_functional_slides(slides_text, llm)
    logger.info(f"[归纳] 功能页: {functional_map}, 内容页: {len(content_indices)} 张")

    # 3. 结构聚类
    clusters = _cluster_content_slides(prs, content_indices)
    logger.info(f"[归纳] 聚类结果: {len(clusters)} 个布局组")

    # 4. Schema 提取
    result = {}
    for cluster in clusters:
        if not cluster:
            continue
        schema = _extract_schema(cluster, prs, llm)
        result.update(schema)

    # 5. 添加功能页
    for func_type, slide_numbers in functional_map.items():
        if not slide_numbers:
            continue
        # 用第一张功能页作为代表
        rep_idx = slide_numbers[0] - 1  # 转为 0-based
        if rep_idx < 0 or rep_idx >= total_slides:
            continue
        slide = prs.slides[rep_idx]
        elements = []
        for i, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                preview = _get_shape_preview(shape, max_len=200)
                if preview:
                    # 为功能页推断元素名称
                    name = _infer_functional_element_name(func_type, i, shape, slide)
                    elements.append({
                        "name": name,
                        "type": "text",
                        "data": [preview],
                    })

        result[func_type] = {
            "template_id": slide_numbers[0],  # 1-based
            "slides": slide_numbers,
            "elements": elements,
        }

    # 6. 添加元数据
    functional_keys = list(functional_map.keys())
    result["functional_keys"] = functional_keys

    # 简单语言检测
    all_text = " ".join(s["title"] + s["body"] for s in slides_text)
    cn_chars = sum(1 for c in all_text if '一' <= c <= '鿿')
    lid = "zh" if cn_chars > len(all_text) * 0.1 else "en"
    result["language"] = {"lid": lid}

    logger.info(f"[归纳] 完成: {len(result) - 2} 个布局")
    return result


def _infer_functional_element_name(func_type: str, shape_idx: int, shape, slide) -> str:
    """为功能页的形状推断语义名称"""
    preview = _get_shape_preview(shape, max_len=50)
    preview_lower = preview.lower()

    if func_type == "opening":
        # 封面页：最高最大的文字通常是主标题
        if shape.top and shape.top < Emu(2500000):
            return "main title"
        elif any(kw in preview_lower for kw in ["by ", "author", "presenter", "日期", "date"]):
            return "presenter"
        else:
            return "section subtitle"
    elif func_type == "table of contents":
        if shape_idx == 0:
            return "main title"
        return "section list"
    elif func_type == "section outline":
        if shape_idx == 0:
            return "main title"
        return "main body"
    elif func_type == "ending":
        if shape_idx == 0:
            return "main title"
        return "main paragraph"
    return f"text_{shape_idx}"


# ==================== CLI 入口 ====================

if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("Usage: python pptx_inductor.py <pptx_path> [output_path]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None

    logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(name)s] %(levelname)s: %(message)s")

    result = induct_template(pptx_path)

    if output_path:
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(result, f, ensure_ascii=False, indent=2)
        print(f"归纳结果已保存到: {output_path}")
    else:
        print(json.dumps(result, ensure_ascii=False, indent=2))
