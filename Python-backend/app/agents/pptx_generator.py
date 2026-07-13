"""
PPT 生成智能体 v9 — 双 Agent + 视觉多样化 + 智能重写 + 富文本 + 双编码
流程: 样式分析 → 内容手稿 → 布局选择(2a) → 内容填充(2b) → 双编码解析 → 功能页插入 → 智能重写 → python-pptx

参照 PPTAgent/DeepPresenter + skills/pptx 规范:
- 双 Agent 机制：布局选择与内容生成分离（参照 PPTAgent layout_selector + editor）
- 视觉多样化：每种布局有 2-3 种渲染变体，通过 slide_index 选择避免千篇一律
- 智能重写：超长文本调用 LLM 压缩而非硬截断（参照 PPTAgent length_rewrite）
- 富文本渲染：Markdown 粗体/斜体/代码格式保留为 PPTX 样式
- 双编码：content/two_column 布局同时输出段落形式和要点形式
- 功能页面自动插入：TOC 和 Section 分隔页自动补齐
- 内容感知样式：LLM 根据主题/领域/受众 + 用户画像动态选择配色方案
- 布局元数据：每种布局的元素有 suggested_characters 限制
- 候选布局随机打乱避免位置偏差
- LLM 调用 3 次重试
- 原生 PPTX 元素：文字框、形状，可编辑
"""
import io
import uuid
import re
import json
import random
import logging
from pathlib import Path
from typing import Dict, Any, List, Tuple
from string import Template

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from app.agents.llm_factory import get_resource_type_generator_llm, MIMOClient
from app.agents.search_utils import search_tavily, download_image
from app.utils.token_recorder import record_from_mimo

logger = logging.getLogger("agents.pptx_generator")

PPTX_DIR = Path(__file__).resolve().parent.parent.parent / "temp" / "pptx"
PPTX_DIR.mkdir(parents=True, exist_ok=True)

# ─── 默认主题色彩（深蓝 + 紫色）───
_DEFAULT_THEME = {
    "pri": "#1B2A4A",
    "pri2": "#2d3f66",
    "acc": "#7C3AED",
    "acc2": "#A78BFA",
    "acc3": "#EDE9FE",
    "bg": "#FFFFFF",
    "bg2": "#F8F7FF",
    "bg3": "#F1F5F9",
    "txt": "#1E293B",
    "txt2": "#64748B",
    "txt3": "#94A3B8",
    "bdr": "#E2E8F0",
}

# 当前活跃主题（运行时动态更新）
_C = {}


def _hex2rgb(hex_str: str) -> RGBColor:
    """'#RRGGBB' → RGBColor"""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _apply_theme(theme_dict: dict):
    """将 hex 颜色字典应用到全局 _C 主题，容错处理非法值"""
    global _C
    _C = {}
    for k, v in theme_dict.items():
        if k == "white":
            continue
        if isinstance(v, str) and re.match(r'^#[0-9a-fA-F]{6}$', v):
            _C[k] = _hex2rgb(v)
        else:
            # 回退到默认值
            default = _DEFAULT_THEME.get(k)
            if default:
                _C[k] = _hex2rgb(default)
            logger.warning(f"[PPT] 主题字段 {k}={v} 非合法 hex，使用默认值")
    _C["white"] = RGBColor(0xFF, 0xFF, 0xFF)


# 初始化默认主题
_apply_theme(_DEFAULT_THEME)

FONT_CN = "Microsoft YaHei"
FONT_MONO = "Cascadia Code"

# 幻灯片尺寸 (16:9)
SW = Inches(13.333)
SH = Inches(7.5)


def _strip_md(text: str) -> str:
    """去除 Markdown 格式标记，返回纯文本"""
    if not text:
        return ""
    t = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    t = re.sub(r'\*(.*?)\*', r'\1', t)
    t = re.sub(r'`(.*?)`', r'\1', t)
    return t


def _add_rich_runs(paragraph, text: str, font_size: int = 16, color=None, font_name=None):
    """解析 Markdown 格式并向 paragraph 添加带样式的 runs。
    支持: **粗体**, *斜体*, `代码`
    """
    if not text:
        return
    base_color = color or _C["txt"]
    base_font = font_name or FONT_CN
    # 正则匹配: **bold**, *italic*, `code`, 普通文本
    pattern = re.compile(r'\*\*(.*?)\*\*|\*(.*?)\*|`(.*?)`|([^*`]+)')
    for m in pattern.finditer(text):
        bold_text, italic_text, code_text, plain_text = m.groups()
        run = paragraph.add_run()
        if bold_text is not None:
            run.text = bold_text
            run.font.bold = True
            run.font.color.rgb = base_color
        elif italic_text is not None:
            run.text = italic_text
            run.font.italic = True
            run.font.color.rgb = base_color
        elif code_text is not None:
            run.text = code_text
            run.font.name = FONT_MONO
            run.font.size = Pt(font_size - 1)
            run.font.color.rgb = _C["acc"]
        else:
            run.text = plain_text
            run.font.color.rgb = base_color
        run.font.size = Pt(font_size)
        run.font.name = base_font


# ─── 原生 PPTX 构建工具 ───

def _add_rect(slide, left, top, width, height, fill=None, line_color=None, line_width=None):
    """添加矩形/圆形装饰"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()  # 无边框
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.fill.solid()
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt(1)
    return shape


def _add_circle(slide, left, top, size, fill):
    """添加圆形装饰"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def _add_image_to_slide(slide, image_data: io.BytesIO, left, top, max_width, max_height):
    """将 BytesIO 图片插入幻灯片，自动缩放保持宽高比"""
    from PIL import Image
    try:
        image_data.seek(0)
        img = Image.open(image_data)
        img_w, img_h = img.size
        if img_w <= 0 or img_h <= 0:
            return False

        # 计算缩放比例
        ratio_w = max_width / img_w
        ratio_h = max_height / img_h
        ratio = min(ratio_w, ratio_h)
        final_w = int(img_w * ratio)
        final_h = int(img_h * ratio)

        # 居中放置在给定区域内
        offset_x = (max_width - final_w) // 2
        offset_y = (max_height - final_h) // 2

        # python-pptx 不支持 WEBP 等格式，统一转为 PNG
        if img.format and img.format.upper() not in ("BMP", "GIF", "JPEG", "PNG", "TIFF", "WMF"):
            buf = io.BytesIO()
            img.convert("RGB").save(buf, format="PNG")
            buf.seek(0)
            image_data = buf

        image_data.seek(0)
        slide.shapes.add_picture(image_data, left + offset_x, top + offset_y, final_w, final_h)
        return True
    except Exception as e:
        logger.warning(f"[PPT] 图片插入失败: {e}")
        return False


def _color_hex(c):
    """RGBColor -> 'RRGGBB' hex string"""
    return f'{c[0]:02X}{c[1]:02X}{c[2]:02X}'


def _alpha_val(percent):
    """百分比不透明度 -> OOXML alpha 值 (100000=100%)"""
    return int(percent * 1000)


def _add_gradient_circle(slide, left, top, size, color1, color2, alpha1=0, alpha2=0):
    """添加渐变圆形（alpha: 0-100 百分比不透明度）"""
    from pptx.oxml.ns import qn
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.line.fill.background()
    sp_pr = shape._element.find(qn('p:spPr'))
    for old_fill in sp_pr.findall(qn('a:solidFill')):
        sp_pr.remove(old_fill)
    for old_fill in sp_pr.findall(qn('a:gradFill')):
        sp_pr.remove(old_fill)
    grad = sp_pr.makeelement(qn('a:gradFill'), {})
    gs_lst = grad.makeelement(qn('a:gsLst'), {})
    gs1 = gs_lst.makeelement(qn('a:gs'), {'pos': '0'})
    srgb1 = gs1.makeelement(qn('a:srgbClr'), {'val': _color_hex(color1)})
    if alpha1 > 0:
        srgb1.append(srgb1.makeelement(qn('a:alpha'), {'val': str(_alpha_val(alpha1))}))
    gs1.append(srgb1)
    gs_lst.append(gs1)
    gs2 = gs_lst.makeelement(qn('a:gs'), {'pos': '100000'})
    srgb2 = gs2.makeelement(qn('a:srgbClr'), {'val': _color_hex(color2)})
    if alpha2 > 0:
        srgb2.append(srgb2.makeelement(qn('a:alpha'), {'val': str(_alpha_val(alpha2))}))
    gs2.append(srgb2)
    gs_lst.append(gs2)
    grad.append(gs_lst)
    lin = grad.makeelement(qn('a:lin'), {'ang': '5400000', 'scaled': '1'})
    grad.append(lin)
    sp_pr.append(grad)
    return shape


def _add_gradient_rect(slide, left, top, width, height, color1, color2, angle=0, alpha1=0, alpha2=0):
    """添加渐变矩形（alpha: 0-100 百分比不透明度）"""
    from pptx.oxml.ns import qn
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    sp_pr = shape._element.find(qn('p:spPr'))
    for old_fill in sp_pr.findall(qn('a:solidFill')):
        sp_pr.remove(old_fill)
    for old_fill in sp_pr.findall(qn('a:gradFill')):
        sp_pr.remove(old_fill)
    grad = sp_pr.makeelement(qn('a:gradFill'), {})
    gs_lst = grad.makeelement(qn('a:gsLst'), {})
    gs1 = gs_lst.makeelement(qn('a:gs'), {'pos': '0'})
    srgb1 = gs1.makeelement(qn('a:srgbClr'), {'val': _color_hex(color1)})
    if alpha1 > 0:
        srgb1.append(srgb1.makeelement(qn('a:alpha'), {'val': str(_alpha_val(alpha1))}))
    gs1.append(srgb1)
    gs_lst.append(gs1)
    gs2 = gs_lst.makeelement(qn('a:gs'), {'pos': '100000'})
    srgb2 = gs2.makeelement(qn('a:srgbClr'), {'val': _color_hex(color2)})
    if alpha2 > 0:
        srgb2.append(srgb2.makeelement(qn('a:alpha'), {'val': str(_alpha_val(alpha2))}))
    gs2.append(srgb2)
    gs_lst.append(gs2)
    grad.append(gs_lst)
    lin = grad.makeelement(qn('a:lin'), {'ang': str(angle * 60000), 'scaled': '1'})
    grad.append(lin)
    sp_pr.append(grad)
    return shape


def _add_scattered_shapes(slide, region="right", count=5, palette=None, seed=None):
    """散点装饰：在指定区域随机放置多个小几何形状

    Args:
        region: "right" / "left" / "corners" / "edges"
        count: 形状数量
        palette: 颜色列表，默认 [_C["acc"], _C["acc2"], _C["acc3"]]
        seed: 随机种子（None 则每次不同）
    """
    if seed is not None:
        rng = random.Random(seed)
    else:
        rng = random.Random()

    palette = palette or [_C["acc"], _C["acc2"], _C["acc3"]]
    shape_types = [MSO_SHAPE.OVAL, MSO_SHAPE.RECTANGLE, MSO_SHAPE.ROUNDED_RECTANGLE]

    # 根据 region 确定可用区域 (left, top 范围，单位: inches)
    regions = {
        "right":    (8.0, 12.0, -0.5, 7.5),
        "left":     (-1.0, 3.0, -0.5, 7.5),
        "corners":  [(9.5, 12.5, -0.5, 1.5), (9.5, 12.5, 5.5, 7.5), (-1.0, 1.5, -0.5, 1.5), (-1.0, 1.5, 5.5, 7.5)],
        "edges":    (0.0, 13.33, -0.5, 7.5),
    }

    area = regions.get(region, regions["right"])

    for _ in range(count):
        # 选择区域
        if isinstance(area, list):
            a = rng.choice(area)
        else:
            a = area

        x = rng.uniform(a[0], a[1])
        y = rng.uniform(a[2], a[3])
        size = rng.uniform(0.15, 0.6)
        color = rng.choice(palette)
        shape_type = rng.choice(shape_types)

        shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(size), Inches(size))
        shape.line.fill.background()
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        # 设置透明度（通过 XML）
        from pptx.oxml.ns import qn
        sp_pr = shape._element.find(qn('p:spPr'))
        solid_fill = sp_pr.find(qn('a:solidFill'))
        if solid_fill is not None:
            srgb = solid_fill.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha_val = rng.randint(15, 40)  # 15%-40% 不透明度
                srgb.append(srgb.makeelement(qn('a:alpha'), {'val': str(_alpha_val(alpha_val))}))


def _add_diagonal_block(slide, color, alpha=20, direction="left_to_right"):
    """对角线色块装饰：一个大三角形斜切页面

    Args:
        color: 填充颜色
        alpha: 不透明度 (0-100)
        direction: "left_to_right" 或 "right_to_left"
    """
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml

    # 用自定义自由形状（三角形）实现对角线
    if direction == "left_to_right":
        # 左下角到右上角的三角形
        points = [
            (Inches(0), Inches(7.5)),      # 左下
            (Inches(13.33), Inches(7.5)),   # 右下
            (Inches(13.33), Inches(3.5)),   # 右上
        ]
    else:
        # 右下角到左上角的三角形
        points = [
            (Inches(0), Inches(7.5)),      # 左下
            (Inches(13.33), Inches(7.5)),   # 右下
            (Inches(0), Inches(3.5)),       # 左上
        ]

    # 用自由形状绘制三角形
    freeform = slide.shapes.build_freeform(points[0][0], points[0][1])
    freeform.add_line_segments(points[1:])
    shape = freeform.convert_to_shape()

    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

    # 设置透明度
    sp_pr = shape._element.find(qn('p:spPr'))
    solid_fill = sp_pr.find(qn('a:solidFill'))
    if solid_fill is not None:
        srgb = solid_fill.find(qn('a:srgbClr'))
        if srgb is not None:
            srgb.append(srgb.makeelement(qn('a:alpha'), {'val': str(_alpha_val(alpha))}))


def _add_textbox(slide, left, top, width, height, text="", font_size=18,
                 color=None, bold=False, align=PP_ALIGN.LEFT, font_name=None,
                 anchor=MSO_ANCHOR.TOP, line_spacing=1.5):
    """添加文本框"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    txbox.text_frame.word_wrap = True
    txbox.text_frame.auto_size = None

    p = txbox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color or _C["txt"]
    p.font.bold = bold
    p.font.name = font_name or FONT_CN
    p.alignment = align
    p.line_spacing = Pt(int(font_size * line_spacing))

    txbox.text_frame.paragraphs[0].space_before = Pt(0)
    txbox.text_frame.paragraphs[0].space_after = Pt(0)

    # 垂直对齐
    try:
        txbox.text_frame.paragraphs[0]
        body_props = txbox.text_frame._txBody.find('{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        if body_props is not None:
            anchor_map = {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}
            body_props.set("anchor", anchor_map.get(anchor, "t"))
    except Exception:
        pass

    return txbox


def _add_bullet_list(slide, left, top, width, height, items: list,
                     font_size=16, color=None, numbered=False):
    """添加带序号/符号的要点列表（支持 Markdown 富文本）"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, item in enumerate(items):
        prefix = f"{i+1}. " if numbered else "●  "

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.line_spacing = Pt(int(font_size * 1.6))
        p.space_before = Pt(6)
        p.space_after = Pt(6)

        # 序号/符号 run
        run_prefix = p.add_run()
        run_prefix.text = prefix
        run_prefix.font.size = Pt(font_size - 1)
        run_prefix.font.color.rgb = _C["acc"]
        run_prefix.font.bold = True
        run_prefix.font.name = FONT_CN

        # 富文本正文 runs（解析 Markdown 格式）
        _add_rich_runs(p, str(item), font_size=font_size, color=color)

    return txbox


def _add_title_bar(slide, title: str, subtitle: str = ""):
    """通用标题栏：左侧紫色竖条 + 标题（无 accent line，避免 AI 感）"""
    # 左侧紫色竖条
    _add_rect(slide, Inches(0.7), Inches(0.5), Inches(0.08), Inches(0.5), fill=_C["acc"])

    # 标题 — 36pt 粗体，符合 skill 规范
    _add_textbox(slide, Inches(1.0), Inches(0.42), Inches(10), Inches(0.65),
                 text=title, font_size=36, color=_C["pri"], bold=True)


def _slide_bg(slide, color=None):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color or _C["bg"]


# ─── 各布局的原生 PPTX 构建函数 ───

def _pptx_title(prs, title: str, subtitle: str):
    """封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(slide, _C["pri"])

    # 对角线色块（右下角三角形，20% 不透明度）
    _add_diagonal_block(slide, _C["acc"], alpha=15, direction="left_to_right")

    # 渐变大圆（右上角）
    _add_gradient_circle(slide, Inches(9.5), Inches(-1.5), Inches(5),
                         _C["pri2"], _C["acc"], alpha1=60, alpha2=30)

    # 散点小装饰（右半区）
    _add_scattered_shapes(slide, region="right", count=6,
                          palette=[_C["acc2"], _C["acc3"], _C["pri2"]])

    # 标题 — 44pt 粗体
    _add_textbox(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1.4),
                 text=title, font_size=44, color=_C["white"], bold=True, line_spacing=1.3)

    # 副标题
    if subtitle:
        _add_textbox(slide, Inches(1.5), Inches(4.4), Inches(9), Inches(0.6),
                     text=subtitle, font_size=18, color=RGBColor(0xA0, 0xAE, 0xC0))


def _pptx_toc(prs, data: dict, variant: int = 0):
    """目录页 — 编号卡片"""
    items = data.get("items", [])
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(slide)
    _add_title_bar(slide, "目录")

    cols = min(len(items[:6]), 3)
    card_w = Inches(3.5)
    card_h = Inches(1.3)
    gap = Inches(0.35)
    start_x = Inches(1.0)
    start_y = Inches(1.6)

    for i, item in enumerate(items[:6]):
        r, c = divmod(i, cols)
        x = start_x + c * (card_w + gap)
        y = start_y + r * (card_h + gap)
        _add_rect(slide, x, y, card_w, card_h, fill=_C["bg2"])
        num_size = Inches(0.45)
        _add_circle(slide, x + Inches(0.2), y + Inches(0.2), num_size, _C["acc"])
        _add_textbox(slide, x + Inches(0.2), y + Inches(0.22), num_size, num_size,
                     text=str(i + 1), font_size=15, color=_C["white"], bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_textbox(slide, x + Inches(0.8), y + Inches(0.25), card_w - Inches(1.0), Inches(0.8),
                     text=_strip_md(str(item)), font_size=15, color=_C["txt"], line_spacing=1.4)


def _pptx_content_v0(prs, slide, title: str, bullets: list):
    """content 变体0：经典 — 左侧紫色竖条 + 编号要点"""
    _add_title_bar(slide, title)
    _add_bullet_list(slide, Inches(1.0), Inches(1.5), Inches(11), Inches(5.5),
                     bullets[:5], font_size=17, numbered=True)
    # 右上角渐变圆装饰
    _add_gradient_circle(slide, Inches(11), Inches(-0.3), Inches(2.5),
                         _C["acc"], _C["acc3"], alpha1=25, alpha2=10)


def _pptx_content_v1(prs, slide, title: str, bullets: list):
    """content 变体1：卡片式 — 标题横幅 + 每个要点在浅色卡片中"""
    # 顶部标题横幅
    _add_rect(slide, Inches(0), Inches(0), SW, Inches(1.2), fill=_C["pri"])
    _add_textbox(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
                 text=title, font_size=32, color=_C["white"], bold=True)

    # 右下角散点装饰（3 个小形状）
    _add_scattered_shapes(slide, region="right", count=3,
                          palette=[_C["acc3"], _C["bg3"]])

    # 每个要点一个卡片
    card_h = Inches(0.85)
    gap = Inches(0.15)
    start_y = Inches(1.6)
    for i, item in enumerate(bullets[:5]):
        y = start_y + i * (card_h + gap)
        _add_rect(slide, Inches(0.8), y, Inches(11.5), card_h, fill=_C["bg2"])
        # 序号圆
        _add_circle(slide, Inches(1.0), y + Inches(0.15), Inches(0.5), _C["acc"])
        _add_textbox(slide, Inches(1.0), y + Inches(0.17), Inches(0.5), Inches(0.5),
                     text=str(i + 1), font_size=16, color=_C["white"], bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 文字
        _add_textbox(slide, Inches(1.7), y + Inches(0.15), Inches(10.2), Inches(0.55),
                     text=_strip_md(str(item)), font_size=15, color=_C["txt"],
                     anchor=MSO_ANCHOR.MIDDLE)


def _pptx_content_v2(prs, slide, title: str, bullets: list):
    """content 变体2：大数字 — 左侧大号序号 + 右侧标题和描述"""
    _add_title_bar(slide, title)

    # 左下角渐变装饰条
    _add_gradient_rect(slide, Inches(0), Inches(6.5), Inches(3), Inches(1),
                       _C["acc3"], _C["bg"], angle=0, alpha1=30, alpha2=0)

    for i, item in enumerate(bullets[:5]):
        y = Inches(1.6) + i * Inches(1.1)
        # 大号序号
        _add_textbox(slide, Inches(0.8), y, Inches(1.2), Inches(0.9),
                     text=f"0{i+1}", font_size=36, color=_C["acc2"], bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 分隔竖线
        _add_rect(slide, Inches(2.1), y + Inches(0.1), Pt(2), Inches(0.7), fill=_C["acc3"])
        # 要点文字
        _add_textbox(slide, Inches(2.4), y + Inches(0.05), Inches(10), Inches(0.8),
                     text=_strip_md(str(item)), font_size=16, color=_C["txt"],
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.5)


def _pptx_content(prs, data: dict, variant: int = 0):
    """标准内容页 — 有图片时用图文混排变体，否则走原有 3 种变体"""
    if data.get("_image_data"):
        _pptx_content_with_image(prs, data)
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(slide)
    title = data.get("title", "")
    bullets = data.get("bullets", [])
    [_pptx_content_v0, _pptx_content_v1, _pptx_content_v2][variant % 3](prs, slide, title, bullets)


def _pptx_content_with_image(prs, data: dict):
    """图文混排内容页：左侧文字 + 右侧图片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(slide)
    title = data.get("title", "")
    bullets = data.get("bullets", [])
    image_data = data.get("_image_data")

    # 标题栏：全宽深色横幅
    _add_rect(slide, Inches(0), Inches(0), SW, Inches(1.2), fill=_C["pri"])
    _add_textbox(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
                 text=title, font_size=32, color=_C["white"], bold=True)

    # 左侧文字区（65% 宽度）
    text_area_w = Inches(7.2)
    start_y = Inches(1.6)
    for i, item in enumerate(bullets[:3]):
        y = start_y + i * Inches(1.6)
        # 大号序号
        _add_textbox(slide, Inches(0.8), y, Inches(1.0), Inches(0.8),
                     text=f"0{i+1}", font_size=30, color=_C["acc2"], bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 分隔竖线
        _add_rect(slide, Inches(1.9), y + Inches(0.1), Pt(2), Inches(0.6), fill=_C["acc3"])
        # 要点文字（用 paragraph 形式，完整句子）
        _add_textbox(slide, Inches(2.2), y + Inches(0.05), Inches(5.0), Inches(0.7),
                     text=_strip_md(str(item)), font_size=15, color=_C["txt"],
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.5)

    # 右侧图片区
    img_left = Inches(7.8)
    img_top = Inches(1.6)
    img_max_w = Inches(4.8)
    img_max_h = Inches(5.2)

    # 图片细边框背景（比图片区域略大）
    frame_pad = Inches(0.08)
    _add_rect(slide, img_left - frame_pad, img_top - frame_pad,
              img_max_w + frame_pad * 2, img_max_h + frame_pad * 2,
              fill=None, line_color=_C.get("bdr", _hex2rgb("#E2E8F0")), line_width=Pt(0.75))

    _add_image_to_slide(slide, image_data, img_left, img_top, img_max_w, img_max_h)


def _pptx_two_column_v0(prs, slide, title: str, left: list, right: list):
    """two_column 变体0：卡片对比 — 左右两个浅色卡片"""
    _add_title_bar(slide, title)
    col_w = Inches(5.5)
    col_h = Inches(5.2)
    y = Inches(1.5)

    _add_rect(slide, Inches(1.0), y, col_w, col_h, fill=_C["bg2"])
    _add_textbox(slide, Inches(1.3), y + Inches(0.2), Inches(2), Inches(0.4),
                 text="对比 A", font_size=13, color=_C["acc"], bold=True)
    _add_bullet_list(slide, Inches(1.3), y + Inches(0.7), col_w - Inches(0.6), col_h - Inches(1.0),
                     left[:4], font_size=15, color=_C["txt"])

    _add_rect(slide, Inches(7.0), y, col_w, col_h, fill=_C["bg2"])
    _add_textbox(slide, Inches(7.3), y + Inches(0.2), Inches(2), Inches(0.4),
                 text="对比 B", font_size=13, color=_C["pri"], bold=True)
    _add_bullet_list(slide, Inches(7.3), y + Inches(0.7), col_w - Inches(0.6), col_h - Inches(1.0),
                     right[:4], font_size=15, color=_C["txt"])


def _pptx_two_column_v1(prs, slide, title: str, left: list, right: list):
    """two_column 变体1：分隔线 — 中间竖线 + 左右标题色带"""
    _add_title_bar(slide, title)
    y = Inches(1.5)
    col_w = Inches(5.5)

    # 左栏标题色带
    _add_rect(slide, Inches(0.8), y, col_w, Inches(0.5), fill=_C["acc"])
    _add_textbox(slide, Inches(1.0), y + Inches(0.05), Inches(3), Inches(0.4),
                 text="对比 A", font_size=14, color=_C["white"], bold=True)
    # 左栏内容
    _add_bullet_list(slide, Inches(1.0), y + Inches(0.7), col_w - Inches(0.4), Inches(4.5),
                     left[:4], font_size=15, color=_C["txt"])

    # 中间竖线
    _add_rect(slide, Inches(6.6), y, Pt(2), Inches(5.2), fill=_C["bdr"])

    # 右栏标题色带
    _add_rect(slide, Inches(6.8), y, col_w, Inches(0.5), fill=_C["pri"])
    _add_textbox(slide, Inches(7.0), y + Inches(0.05), Inches(3), Inches(0.4),
                 text="对比 B", font_size=14, color=_C["white"], bold=True)
    # 右栏内容
    _add_bullet_list(slide, Inches(7.0), y + Inches(0.7), col_w - Inches(0.4), Inches(4.5),
                     right[:4], font_size=15, color=_C["txt"])


def _pptx_two_column(prs, data: dict, variant: int = 0):
    """双栏对比页 — 有图片时图片替代右栏，否则走原有 2 种变体"""
    image_data = data.get("_image_data")
    if image_data:
        _pptx_two_column_with_image(prs, data)
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(slide)
    [_pptx_two_column_v0, _pptx_two_column_v1][variant % 2](
        prs, slide, data.get("title", ""), data.get("left", []), data.get("right", []))


def _pptx_two_column_with_image(prs, data: dict):
    """图文混排双栏页：左栏文字 + 右栏图片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(slide)
    title = data.get("title", "")
    left = data.get("left", [])
    image_data = data.get("_image_data")

    # 标题栏
    _add_title_bar(slide, title)

    y = Inches(1.5)
    col_w = Inches(5.5)

    # 左栏文字
    _add_rect(slide, Inches(0.8), y, col_w, Inches(5.2), fill=_C["bg2"])
    _add_textbox(slide, Inches(1.0), y + Inches(0.2), Inches(2), Inches(0.4),
                 text="要点", font_size=13, color=_C["acc"], bold=True)
    _add_bullet_list(slide, Inches(1.0), y + Inches(0.7), col_w - Inches(0.4), Inches(4.3),
                     left[:4], font_size=15, color=_C["txt"])

    # 右栏图片
    img_left = Inches(6.8)
    img_top = y
    img_max_w = Inches(5.8)
    img_max_h = Inches(5.2)

    frame_pad = Inches(0.08)
    _add_rect(slide, img_left - frame_pad, img_top - frame_pad,
              img_max_w + frame_pad * 2, img_max_h + frame_pad * 2,
              fill=None, line_color=_C.get("bdr", _hex2rgb("#E2E8F0")), line_width=Pt(0.75))

    _add_image_to_slide(slide, image_data, img_left, img_top, img_max_w, img_max_h)


def _pptx_cards_v0(prs, slide, title: str, cards: list):
    """cards 变体0：标准竖排 — 大数字 + 标签 + 描述"""
    _add_title_bar(slide, title)
    n = min(len(cards), 4)
    card_w = Inches(2.7)
    card_h = Inches(4.5)
    gap = Inches(0.35)
    total_w = n * card_w + (n - 1) * gap
    start_x = (SW - total_w) // 2
    y = Inches(1.8)

    for i, c in enumerate(cards[:4]):
        x = start_x + i * (card_w + gap)
        num, label, desc = c.get("number", ""), c.get("label", ""), c.get("desc", "")
        _add_rect(slide, x, y, card_w, card_h, fill=_C["white"])
        inner_x = x + Inches(0.2)
        cur_y = y + Inches(0.4)
        if num:
            _add_textbox(slide, inner_x, cur_y, card_w - Inches(0.4), Inches(0.8),
                         text=str(num), font_size=38, color=_C["acc"], bold=True,
                         align=PP_ALIGN.CENTER)
            cur_y += Inches(0.9)
        _add_textbox(slide, inner_x, cur_y, card_w - Inches(0.4), Inches(0.5),
                     text=_strip_md(str(label)), font_size=17, color=_C["pri"], bold=True,
                     align=PP_ALIGN.CENTER)
        cur_y += Inches(0.6)
        if desc:
            _add_textbox(slide, inner_x, cur_y, card_w - Inches(0.4), Inches(2),
                         text=_strip_md(str(desc)), font_size=13, color=_C["txt2"],
                         align=PP_ALIGN.CENTER, line_spacing=1.5)


def _pptx_cards_v1(prs, slide, title: str, cards: list):
    """cards 变体1：横向条纹 — 交替色带 + 数字在左、说明在右"""
    _add_title_bar(slide, title)
    row_h = Inches(1.2)
    gap = Inches(0.2)
    start_y = Inches(1.6)
    stripe_colors = [_C["bg2"], _C["bg3"]]

    for i, c in enumerate(cards[:4]):
        y = start_y + i * (row_h + gap)
        num, label, desc = c.get("number", ""), c.get("label", ""), c.get("desc", "")
        # 交替色带背景
        _add_rect(slide, Inches(0.8), y, Inches(11.5), row_h, fill=stripe_colors[i % 2])
        # 左侧数字
        if num:
            _add_textbox(slide, Inches(1.0), y + Inches(0.1), Inches(2), Inches(1.0),
                         text=str(num), font_size=32, color=_C["acc"], bold=True,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 右侧标签+描述
        _add_textbox(slide, Inches(3.2), y + Inches(0.15), Inches(8.5), Inches(0.4),
                     text=_strip_md(str(label)), font_size=16, color=_C["pri"], bold=True)
        if desc:
            _add_textbox(slide, Inches(3.2), y + Inches(0.55), Inches(8.5), Inches(0.5),
                         text=_strip_md(str(desc)), font_size=13, color=_C["txt2"],
                         line_spacing=1.4)


def _pptx_cards(prs, data: dict, variant: int = 0):
    """数据卡片页 — 2种视觉变体"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(slide)
    [_pptx_cards_v0, _pptx_cards_v1][variant % 2](
        prs, slide, data.get("title", ""), data.get("cards", []))


def _pptx_section(prs, data: dict, variant: int = 0):
    """章节分隔页"""
    title = data.get("title", "")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(slide, _C["bg2"])

    # 渐变大圆 + 散点装饰
    _add_gradient_circle(slide, Inches(9), Inches(-1.5), Inches(4.5),
                         _C["acc3"], _C["acc"], alpha1=50, alpha2=15)
    _add_scattered_shapes(slide, region="corners", count=4,
                          palette=[_C["acc3"], _C["bg3"]])

    _add_textbox(slide, Inches(2), Inches(3.0), Inches(9), Inches(1.5),
                 text=_strip_md(title), font_size=40, color=_C["acc"], bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)


def _pptx_code(prs, data: dict, variant: int = 0):
    """代码/流程页"""
    title = data.get("title", "")
    code = data.get("code", "")
    bullets = data.get("bullets", [])
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(slide)
    _add_title_bar(slide, title)
    code_x, code_y, code_w, code_h = Inches(1.0), Inches(1.5), Inches(6.5), Inches(5.2)
    _add_rect(slide, code_x, code_y, code_w, code_h, fill=RGBColor(0x1E, 0x1E, 0x2E))
    _add_textbox(slide, code_x + Inches(0.3), code_y + Inches(0.3), code_w - Inches(0.6), code_h - Inches(0.6),
                 text=code[:800], font_size=13, color=RGBColor(0xCD, 0xD6, 0xF4),
                 font_name=FONT_MONO, line_spacing=1.7)
    _add_bullet_list(slide, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.5),
                     bullets[:4], font_size=15)


def _pptx_quote_v0(prs, slide, quote: str, source: str, bullets: list):
    """quote variant 0: centered big quote mark"""
    _add_circle(slide, Inches(10), Inches(0.5), Inches(1.5), _C["acc3"])
    _add_textbox(slide, Inches(2.5), Inches(1.2), Inches(1), Inches(1),
                 text="“", font_size=72, color=_C["acc2"], align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(2.5), Inches(2.3), Inches(8), Inches(2),
                 text=_strip_md(quote), font_size=22, color=_C["pri"],
                 align=PP_ALIGN.CENTER, line_spacing=1.6)
    if source:
        _add_textbox(slide, Inches(2.5), Inches(4.3), Inches(8), Inches(0.4),
                     text=_strip_md(source), font_size=14, color=_C["txt2"],
                     align=PP_ALIGN.CENTER)
    if bullets:
        _add_bullet_list(slide, Inches(2.5), Inches(5.3), Inches(8), Inches(1.8),
                         bullets[:3], font_size=14, color=_C["txt2"])


def _pptx_quote_v1(prs, slide, quote: str, source: str, bullets: list):
    """quote variant 1: left color block + right quote text"""
    _add_rect(slide, Inches(0), Inches(0), Inches(4), SH, fill=_C["acc"])
    _add_circle(slide, Inches(2.5), Inches(5.5), Inches(2), _C["acc2"])
    _add_textbox(slide, Inches(1), Inches(1.5), Inches(2), Inches(1.5),
                 text="“", font_size=80, color=_C["white"], align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(5), Inches(2.0), Inches(7.5), Inches(2.5),
                 text=_strip_md(quote), font_size=22, color=_C["pri"],
                 line_spacing=1.6)
    if source:
        _add_textbox(slide, Inches(5), Inches(4.5), Inches(7.5), Inches(0.4),
                     text=f"— {_strip_md(source)}", font_size=14, color=_C["txt2"])
    if bullets:
        _add_bullet_list(slide, Inches(5), Inches(5.2), Inches(7.5), Inches(2),
                         bullets[:3], font_size=14, color=_C["txt2"])


def _pptx_quote(prs, data: dict, variant: int = 0):
    """Quote/Highlight page - 2 visual variants"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(slide, _C["bg2"])
    [_pptx_quote_v0, _pptx_quote_v1][variant % 2](
        prs, slide, data.get("quote", ""), data.get("source", ""), data.get("bullets", []))



def _pptx_timeline_v0(prs, slide, title: str, steps: list):
    """timeline 变体0：水平时间线 — 水平连接线 + 圆形编号"""
    _add_title_bar(slide, title)
    n = min(len(steps), 5)
    step_w = Inches(2.2)
    gap = Inches(0.3)
    total_w = n * step_w + (n - 1) * gap
    start_x = (SW - total_w) // 2
    y_center = Inches(4.0)

    line_y = y_center + Inches(0.25)
    _add_rect(slide, start_x + Inches(0.3), line_y,
              total_w - Inches(0.6), Pt(3), fill=_C["acc3"])

    for i, step in enumerate(steps[:5]):
        x = start_x + i * (step_w + gap)
        label, desc = step.get("label", ""), step.get("desc", "")
        cx = x + step_w // 2 - Inches(0.25)
        _add_circle(slide, cx, y_center, Inches(0.5), _C["acc"])
        _add_textbox(slide, cx, y_center + Inches(0.02), Inches(0.5), Inches(0.5),
                     text=str(i + 1), font_size=18, color=_C["white"], bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_textbox(slide, x, y_center - Inches(0.8), step_w, Inches(0.5),
                     text=_strip_md(str(label)), font_size=15, color=_C["pri"],
                     bold=True, align=PP_ALIGN.CENTER)
        if desc:
            _add_textbox(slide, x, y_center + Inches(0.7), step_w, Inches(1.5),
                         text=_strip_md(str(desc)), font_size=12, color=_C["txt2"],
                         align=PP_ALIGN.CENTER, line_spacing=1.4)


def _pptx_timeline_v1(prs, slide, title: str, steps: list):
    """timeline 变体1：垂直阶梯 — 左侧时间线 + 交替左右的卡片"""
    _add_title_bar(slide, title)
    n = min(len(steps), 5)
    line_x = Inches(2.5)
    start_y = Inches(1.6)
    row_h = Inches(1.1)

    # 垂直连接线
    _add_rect(slide, line_x + Inches(0.22), start_y, Pt(3), row_h * n, fill=_C["acc3"])

    for i, step in enumerate(steps[:5]):
        y = start_y + i * row_h
        label, desc = step.get("label", ""), step.get("desc", "")
        # 编号圆
        _add_circle(slide, line_x, y + Inches(0.15), Inches(0.5), _C["acc"])
        _add_textbox(slide, line_x, y + Inches(0.17), Inches(0.5), Inches(0.5),
                     text=str(i + 1), font_size=16, color=_C["white"], bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 卡片（右侧）
        card_x = Inches(3.3)
        _add_rect(slide, card_x, y + Inches(0.05), Inches(9), Inches(0.85), fill=_C["bg2"])
        _add_textbox(slide, card_x + Inches(0.3), y + Inches(0.1), Inches(2.5), Inches(0.35),
                     text=_strip_md(str(label)), font_size=14, color=_C["acc"], bold=True)
        if desc:
            _add_textbox(slide, card_x + Inches(0.3), y + Inches(0.45), Inches(8.2), Inches(0.35),
                         text=_strip_md(str(desc)), font_size=12, color=_C["txt2"])


def _pptx_timeline(prs, data: dict, variant: int = 0):
    """时间线/流程页 — 2种视觉变体"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(slide)
    [_pptx_timeline_v0, _pptx_timeline_v1][variant % 2](
        prs, slide, data.get("title", ""), data.get("steps", []))


def _pptx_ending(prs, title: str):
    """结尾页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(slide, _C["pri"])

    # 对角线色块（方向与封面相反，形成呼应）
    _add_diagonal_block(slide, _C["acc"], alpha=12, direction="right_to_left")

    # 渐变圆（左下角）
    _add_gradient_circle(slide, Inches(-1.5), Inches(4), Inches(4.5),
                         _C["acc"], _C["pri2"], alpha1=40, alpha2=20)

    # 散点小装饰（左半区 + 角落）
    _add_scattered_shapes(slide, region="left", count=4,
                          palette=[_C["acc2"], _C["acc3"], _C["pri2"]])

    _add_textbox(slide, Inches(2), Inches(2.8), Inches(9), Inches(1.2),
                 text="谢谢观看", font_size=48, color=_C["white"], bold=True,
                 align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(2), Inches(4.2), Inches(9), Inches(0.5),
                 text=title, font_size=16, color=RGBColor(0x90, 0x9E, 0xB0),
                 align=PP_ALIGN.CENTER)


# 布局映射（接受 prs, data, variant 三个参数）
_PPTX_BUILDERS = {
    "toc":        _pptx_toc,
    "content":    _pptx_content,
    "two_column": _pptx_two_column,
    "cards":      _pptx_cards,
    "section":    _pptx_section,
    "code":       _pptx_code,
    "quote":      _pptx_quote,
    "timeline":   _pptx_timeline,
}

# 每种布局的变体数量
_VARIANT_COUNT = {
    "content": 3,
    "two_column": 2,
    "cards": 2,
    "timeline": 2,
    "quote": 2,
    "toc": 1,
    "section": 1,
    "code": 1,
}


# ─── 布局元数据（参照 PPTAgent 的 suggested_characters 机制）───

LAYOUT_META = {
    "toc": {
        "name": "目录页",
        "desc": "列出主要章节，仅用1次放最前",
        "elements": {
            "items": {"type": "text_list", "max_items": 6, "suggested_chars": 20},
        },
    },
    "content": {
        "name": "标准内容页",
        "desc": "带编号要点列表，最常用，占比约40%",
        "elements": {
            "title": {"type": "text", "suggested_chars": 30},
            "bullets": {"type": "text_list", "max_items": 5, "suggested_chars": 80},
        },
    },
    "two_column": {
        "name": "双栏对比页",
        "desc": "适合对比、分类说明",
        "elements": {
            "title": {"type": "text", "suggested_chars": 30},
            "left": {"type": "text_list", "max_items": 4, "suggested_chars": 60},
            "right": {"type": "text_list", "max_items": 4, "suggested_chars": 60},
        },
    },
    "cards": {
        "name": "数据卡片页",
        "desc": "展示关键数据、指标、特征",
        "elements": {
            "title": {"type": "text", "suggested_chars": 30},
            "cards": {"type": "card_list", "max_items": 4, "fields": {
                "number": {"suggested_chars": 10},
                "label": {"suggested_chars": 20},
                "desc": {"suggested_chars": 60},
            }},
        },
    },
    "section": {
        "name": "章节分隔页",
        "desc": "大号标题居中，每个大章节前1次",
        "elements": {
            "title": {"type": "text", "suggested_chars": 25},
        },
    },
    "code": {
        "name": "代码/流程展示页",
        "desc": "左侧代码右侧说明",
        "elements": {
            "title": {"type": "text", "suggested_chars": 30},
            "code": {"type": "text", "suggested_chars": 800},
            "bullets": {"type": "text_list", "max_items": 4, "suggested_chars": 60},
        },
    },
    "quote": {
        "name": "引用/重点页",
        "desc": "突出核心观点或名言",
        "elements": {
            "quote": {"type": "text", "suggested_chars": 120},
            "source": {"type": "text", "suggested_chars": 40},
            "bullets": {"type": "text_list", "max_items": 3, "suggested_chars": 60},
        },
    },
    "timeline": {
        "name": "时间线/流程页",
        "desc": "展示步骤或演进过程",
        "elements": {
            "title": {"type": "text", "suggested_chars": 30},
            "steps": {"type": "step_list", "max_items": 5, "fields": {
                "label": {"suggested_chars": 15},
                "desc": {"suggested_chars": 50},
            }},
        },
    },
}


REWRITE_PROMPT = Template("""你是一个 PPT 内容精简专家。以下是多条超长的 PPT 文本，请将每条精简到指定字数以内。

## 精简原则
1. 保留核心信息和关键词，删除冗余修饰
2. 保持语句通顺，不要出现断句
3. 保留具体数据、数字、专有名词
4. 优先保留结论，删除背景描述

## 待精简文本
$texts_json

## 输出格式
严格输出 JSON 数组，每个元素是精简后的文本，顺序与输入一致:
["精简后的文本1", "精简后的文本2", ...]

请直接输出 JSON:""")


def _validate_and_compress(slides_data: list, llm=None) -> list:
    """参照 PPTAgent 的 length_rewrite：验证文本长度，超长则 LLM 智能重写，失败则截断"""
    # 第一步：收集所有超长文本及其位置
    overlong = []  # [(slide_idx, key, field_key_or_none, original_text, limit)]

    for idx, sd in enumerate(slides_data):
        layout = sd.get("layout", "content")
        meta = LAYOUT_META.get(layout)
        if not meta:
            continue
        for key, spec in meta["elements"].items():
            if key not in sd:
                continue
            if spec["type"] == "text":
                limit = spec["suggested_chars"]
                text = str(sd.get(key, ""))
                if len(text) > limit + 5:
                    overlong.append((idx, key, None, text, limit))
            elif spec["type"] == "text_list":
                limit = spec.get("suggested_chars", 80)
                items = sd.get(key, [])
                if not isinstance(items, list):
                    continue
                sd[key] = items[:spec.get("max_items", 5)]
                for i, item in enumerate(sd[key]):
                    s = str(item)
                    if len(s) > limit + 5:
                        overlong.append((idx, key, i, s, limit))
            elif spec["type"] in ("card_list", "step_list"):
                items = sd.get(key, [])
                if not isinstance(items, list):
                    continue
                sd[key] = items[:spec.get("max_items", 5)]
                for ci, card in enumerate(sd[key]):
                    if not isinstance(card, dict):
                        continue
                    for fkey, fspec in spec.get("fields", {}).items():
                        val = str(card.get(fkey, ""))
                        lim = fspec["suggested_chars"]
                        if len(val) > lim + 5:
                            overlong.append((idx, key, (ci, fkey), val, lim))

    if not overlong:
        return slides_data

    # 第二步：尝试 LLM 批量重写
    if llm:
        try:
            texts_input = [{"text": t[3], "limit": t[4]} for t in overlong]
            prompt = REWRITE_PROMPT.substitute(
                texts_json=json.dumps(texts_input, ensure_ascii=False, indent=1)
            )
            result = llm.chat_json([{"role": "system", "content": prompt}])
            if isinstance(result, list) and len(result) == len(overlong):
                for i, (idx, key, field_key, _, limit) in enumerate(overlong):
                    rewritten = str(result[i])[:limit + 10]  # 留一点余量
                    if field_key is None:
                        slides_data[idx][key] = rewritten
                    elif isinstance(field_key, int):
                        slides_data[idx][key][field_key] = rewritten
                    elif isinstance(field_key, tuple):
                        ci, fkey = field_key
                        slides_data[idx][key][ci][fkey] = rewritten
                logger.info(f"[PPT] LLM 智能重写完成，{len(overlong)} 条文本")
                return slides_data
            else:
                logger.warning(f"[PPT] LLM 重写返回格式不匹配，降级为截断")
        except Exception as e:
            logger.warning(f"[PPT] LLM 重写失败，降级为截断: {e}")

    # 第三步：降级为截断
    for idx, key, field_key, text, limit in overlong:
        truncated = text[:limit] + "..."
        if field_key is None:
            slides_data[idx][key] = truncated
        elif isinstance(field_key, int):
            slides_data[idx][key][field_key] = truncated
        elif isinstance(field_key, tuple):
            ci, fkey = field_key
            slides_data[idx][key][ci][fkey] = truncated
    logger.info(f"[PPT] 截断处理完成，{len(overlong)} 条文本")
    return slides_data


# ─── 样式分析（参照 PPTAgent Design Agent 的内容感知样式选择）───

STYLE_ANALYSIS_PROMPT = Template("""你是一个专业的视觉设计专家。根据 PPT 内容主题和用户画像，选择最合适的配色方案。

## 内容信息
- 标题: $title
- 描述: $description
- 内容摘要: $content_summary

## 用户画像
$user_profile_text

## 配色方案参考（根据主题语义选择最匹配的方案，也可以自定义）

| 方案名 | 主色 | 强调色 | 适用场景 |
|--------|------|--------|---------|
| 深蓝商务 | #1B2A4A | #7C3AED | 商业、管理、金融、通用 |
| 森林自然 | #2C5F2D | #97BC62 | 生物、环境、农业、健康 |
| 珊瑚活力 | #F96167 | #F9E795 | 艺术、创意、儿童教育 |
| 海洋科技 | #065A82 | #1C7293 | 科技、工程、编程、数据 |
| 暖色陶土 | #B85042 | #A7BEAE | 历史、文化、人文社科 |
| 炭灰极简 | #36454F | #212121 | 学术、论文、正式场合 |
| 青绿信任 | #028090 | #02C39A | 医疗、心理、教育 |
| 浆果优雅 | #6D2E46 | #A26769 | 文学、艺术、设计 |
| 鼠尾草 | #84B59F | #69A297 | 哲学、冥想、休闲学习 |
| 樱桃大胆 | #990011 | #FCF6F5 | 竞赛、考试、高强度学习 |

## 输出格式
严格输出 JSON，不要输出其他内容:
{
  "theme_name": "方案名称",
  "pri": "#RRGGBB",
  "pri2": "#RRGGBB",
  "acc": "#RRGGBB",
  "acc2": "#RRGGBB",
  "acc3": "#RRGGBB",
  "bg": "#FFFFFF",
  "bg2": "#RRGGBB",
  "bg3": "#RRGGBB",
  "txt": "#1E293B",
  "txt2": "#64748B",
  "txt3": "#94A3B8",
  "bdr": "#E2E8F0"
}

## 选择原则
1. 根据内容主题的语义选择最匹配的色彩氛围（科技→冷色调、人文→暖色调、自然→绿色系）
2. 如果用户画像中有视觉偏好（视觉型），选择对比度更强的配色
3. 如果用户画像中有学习风格信息，适当调整（如儿童教育→更活泼的颜色）
4. bg/bg3/txt/txt2/txt3/bdr 保持中性色即可，主要调整 pri/pri2/acc/acc2/accent3
5. 确保文字颜色与背景有足够对比度

请直接输出 JSON:""")


def _analyze_style(
    title: str,
    description: str,
    module_content: str,
    user_profile: dict = None,
    user_id: int = 0,
) -> dict:
    """Stage 0: 分析内容主题和用户画像，选择配色方案"""
    llm = get_resource_type_generator_llm()

    # 构建用户画像文本
    profile_text = "暂无用户画像信息"
    if user_profile:
        parts = []
        behavior = user_profile.get("learning_behavior", {})
        if behavior:
            vv = behavior.get("visual_vs_verbal", 0)
            if vv < -0.3:
                parts.append("学习风格：视觉型（偏好图表、色彩丰富的展示）")
            elif vv > 0.3:
                parts.append("学习风格：言语型（偏好文字、逻辑清晰的展示）")
            else:
                parts.append("学习风格：均衡型")
        name = user_profile.get("name") or user_profile.get("nickname", "")
        if name:
            parts.append(f"用户：{name}")
        if parts:
            profile_text = "；".join(parts)

    prompt = STYLE_ANALYSIS_PROMPT.substitute(
        title=title,
        description=description,
        content_summary=module_content[:1500],
        user_profile_text=profile_text,
    )

    for attempt in range(3):
        try:
            result = llm.chat_json([{"role": "system", "content": prompt}])
            # 验证返回的 key 完整性
            required = {"pri", "acc", "bg", "txt"}
            if not required.issubset(result.keys()):
                raise Exception(f"样式分析返回缺少字段: {required - result.keys()}")
            return result
        except Exception as e:
            if attempt == 2:
                logger.warning(f"[PPT] 样式分析失败，使用默认主题: {e}")
                return _DEFAULT_THEME
            logger.warning(f"[PPT] 样式分析重试 {attempt+1}/3: {e}")

    return _DEFAULT_THEME


# ─── LLM Prompts（两阶段）───

MANUSCRIPT_PROMPT = Template("""你是一个专业的教学内容策划专家。你的任务是根据学习材料，撰写一份高质量的 PPT 内容手稿。

## 输出格式
用 Markdown 格式输出，每一页用 `---` 分隔线隔开。格式如下：

# PPT标题
副标题：一句话概括

---
## 目录
- 章节1标题
- 章节2标题
- 章节3标题

---
## 第一页标题
**核心结论**（一句话点明本页要点）

- 要点1：具体且有深度的描述，包含数据、例子或因果分析
- 要点2：...
- 要点3：...
- 要点4：...

---
## 第二页标题
（同上格式，可以是对比、数据、流程等不同内容形式）

---
...更多页面...

## 写作原则

1. **金字塔原则**：每页开头就给出核心结论，要点是支撑论据
2. **信息密度**：每条要点 20~40 个字，要有具体数据、例子、因果关系或深度分析，不要空洞的概括
3. **一页一洞察**：每页只讲一个核心观点，用 3~5 个要点支撑
4. **内容多样**：混合使用以下内容形式：
   - 普通要点列表（最常见）
   - 对比分析（A vs B）
   - 数据/指标（数字和百分比）
   - 流程/步骤（编号步骤）
   - 案例/引言（引用形式）
   - 代码/公式（代码块）
5. **严禁 emoji**
6. **所有文本中文**
7. **总页数 $page_count 页**（不含目录和结尾）

## 学习内容
$module_content

## 模块信息
- 标题: $title
- 描述: $description

请直接输出 Markdown 手稿，不要输出其他内容。第一行必须是 `# PPT标题`。""")


def _build_layouts_text() -> str:
    """构建布局元数据描述文本（随机打乱顺序）"""
    layouts_info = []
    keys = list(LAYOUT_META.keys())
    random.shuffle(keys)
    for k in keys:
        meta = LAYOUT_META[k]
        elems_desc = []
        for ek, es in meta["elements"].items():
            if es["type"] == "text":
                elems_desc.append(f"    {ek}: 文本，建议{es['suggested_chars']}字以内")
            elif es["type"] in ("text_list", "card_list", "step_list"):
                elems_desc.append(f"    {ek}: 列表，最多{es['max_items']}项")
        layouts_info.append(f"- \"{k}\" ({meta['name']}): {meta['desc']}\n" + "\n".join(elems_desc))
    return "\n".join(layouts_info)


IMAGE_QUERY_PROMPT = """你是一个 PPT 配图搜索助手。PPT 有配图会更生动、更专业。请为每个页面生成图片搜索关键词。

## 规则
- 每个 content/two_column 页面都应该配图，除非内容完全不适合（如纯代码、纯引用）
- 根据页面摘要和主题，用**中文**描述要搜索的图片
- 关键词要具体，例如："TCP三次握手流程图" 而不是 "网络"
- 优先搜索：概念图解、流程图、架构图、数据可视化、示意图、实物照片

## 页面信息
$slides_info

## 输出格式
严格输出 JSON 对象，包含 slides 数组，每个元素对应一页：
{{"slides": [{{"index": 0, "image_query": "TCP三次握手流程图"}}, {{"index": 1, "image_query": "拥塞控制算法示意图"}}]}}

请直接输出 JSON:"""


def _image_query_desc(slide: Dict) -> str:
    """返回图片搜索关键词的简短描述"""
    return slide.get("_image_query", "配图")[:30]


def _generate_image_queries(outline_slides: List[Dict], llm) -> Dict[int, str]:
    """批量为 outline 页面生成图片搜索关键词，返回 {slide_index: query}"""
    # 只处理适合配图的布局
    imageable_layouts = ("content", "two_column")
    slides_info = []
    for idx, s in enumerate(outline_slides):
        layout = s.get("layout", "content")
        summary = s.get("summary", "")
        if layout in imageable_layouts:
            slides_info.append(f"第{idx}页 [layout={layout}]: {summary}")

    if not slides_info:
        return {}

    prompt = IMAGE_QUERY_PROMPT.replace("$slides_info", "\n".join(slides_info))
    try:
        result = llm.chat_json([{"role": "system", "content": prompt}])
        logger.info(f"[PPT] 图片查询 LLM 返回: {str(result)[:500]}")
        queries = {}
        items = []
        if isinstance(result, list):
            items = result
        elif isinstance(result, dict):
            items = result.get("slides", result.get("images", []))
            if not items:
                # 兼容单个对象的情况
                items = [result]
        for item in items:
            if not isinstance(item, dict):
                continue
            idx = item.get("index", -1)
            q = item.get("image_query", item.get("query", "")).strip()
            if 0 <= idx < len(outline_slides) and q:
                queries[idx] = q
        logger.info(f"[PPT] 图片查询生成: {len(queries)}/{len(slides_info)} 页需要配图")
        return queries
    except Exception as e:
        logger.warning(f"[PPT] 图片查询生成失败: {e}")
        return {}


def _validate_image_relevance(image_data: io.BytesIO, page_summary: str, user_id: int = 0) -> bool:
    """用 mimo-v2.5 分析图片是否与 PPT 页面内容相关，相关返回 True"""
    import base64
    try:
        image_data.seek(0)
        b64 = base64.b64encode(image_data.read()).decode("utf-8")
        image_data.seek(0)

        llm = MIMOClient(model=MIMOClient.MODEL_STANDARD, temperature=0, max_tokens=256,
                         thinking=None)
        messages = [{
            "role": "user",
            "content": [
                {
                    "type": "text",
                    "text": f"这是一张用于 PPT 配图的图片。PPT 页面的主题是：「{page_summary}」\n\n"
                            f"请判断这张图片是否与主题相关。如果相关（图片内容能辅助说明主题），回答 \"yes\"；"
                            f"如果不相关（图片内容与主题无关、是广告、水印过多、纯文字截图等），回答 \"no\"。\n\n"
                            f"只回答 yes 或 no，不要解释。"
                },
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/jpeg;base64,{b64}"}
                }
            ]
        }]
        result = llm.chat(messages).strip().lower()
        record_from_mimo(llm, user_id, "pptx_image_validate")
        relevant = result.startswith("y")
        if not relevant:
            logger.info(f"[PPT] 图片相关性检查: 不相关 (LLM: {result[:20]})")
        return relevant
    except Exception as e:
        logger.warning(f"[PPT] 图片相关性检查失败: {e}")
        return True  # 检查失败时保留图片


def _search_images_for_slides(outline_slides: List[Dict], image_queries: Dict[int, str],
                               emit_thinking=None) -> List[Dict]:
    """根据搜索关键词并行搜索下载图片，附加到 outline_slides 对应项"""
    if not image_queries:
        return outline_slides

    if emit_thinking:
        emit_thinking(f"正在搜索 {len(image_queries)} 张配图并验证相关性...")

    from concurrent.futures import ThreadPoolExecutor, as_completed
    query_images = {}
    queries_list = list(image_queries.values())
    with ThreadPoolExecutor(max_workers=4) as executor:
        future_map = {executor.submit(search_tavily, q, 3, True): q for q in queries_list}
        for future in as_completed(future_map):
            q = future_map[future]
            try:
                _, images = future.result()
                query_images[q] = images
            except Exception as e:
                logger.warning(f"[PPT] 图片搜索失败 '{q}': {e}")
                query_images[q] = []

    # 下载图片并附加到 outline_slides
    attached = 0
    for idx, query in image_queries.items():
        images = query_images.get(query, [])
        if not images:
            continue
        page_summary = outline_slides[idx].get("summary", query)
        for img in images[:3]:
            url = img.get("url", "")
            if not url:
                continue
            image_data = download_image(url, timeout=10, max_size_mb=5)
            if not image_data:
                continue
            # 用 mimo-v2.5 检查图片是否与内容相关
            if not _validate_image_relevance(image_data, page_summary):
                logger.info(f"[PPT] 第{idx+1}页: 跳过不相关配图")
                continue
            outline_slides[idx]["_image_data"] = image_data
            outline_slides[idx]["_image_query"] = query
            attached += 1
            logger.info(f"[PPT] 第{idx+1}页: 已下载配图 ({len(image_data.getvalue()) // 1024}KB)")
            break

    logger.info(f"[PPT] 图片搜索完成: {attached}/{len(image_queries)} 张配图")
    return outline_slides


OUTLINE_PROMPT = Template("""你是一个专业的 PPT 布局设计师。你的任务是根据内容手稿，为每一页选择最合适的布局类型。

## 输入手稿
---
$manuscript
---

## 可用布局
$layouts_text

## 输出格式
严格输出 JSON，不要输出其他内容:
{{
  "title": "PPT 主标题",
  "subtitle": "副标题",
  "slides": [
    {{"layout": "content", "summary": "本页要讲的核心内容，一句话概括"}},
    {{"layout": "two_column", "summary": "对比什么内容"}},
    {{"layout": "cards", "summary": "展示哪些关键数据"}},
    {{"layout": "section", "summary": "章节标题"}},
    {{"layout": "timeline", "summary": "展示什么流程"}},
    {{"layout": "quote", "summary": "引用什么观点"}}
  ]
}}

## 布局选择规则
1. **内容→布局匹配**：
   - 普通要点/知识点 → content
   - 对比分析(A vs B) → two_column
   - 关键数据/指标/特征 → cards
   - 流程/步骤/演进 → timeline
   - 核心观点/名言/引言 → quote
   - 代码/公式/伪代码 → code
2. **功能布局位置**：toc 仅1次放最前，section 每大章节前1次
3. **严禁连续3页相同布局** — 这是最重要的规则
4. **布局多样性**：整份PPT至少使用3种不同的布局类型
5. **严禁 emoji，所有文本中文**
6. 总页数控制在 8~15 页

请直接输出 JSON:""")


CONTENT_PROMPT = Template("""你是一个专业的 PPT 内容编辑。你的任务是根据手稿原文，为已确定的布局大纲填充具体内容。

## 布局大纲
$outline_json

## 原始手稿
---
$manuscript
---

## 布局元素规格
$layouts_text
$image_info

## 输出格式
严格输出 JSON，不要输出其他内容。对于 content 和 two_column 布局，要点需要提供双编码（paragraph + bullet 两种形式）:
{{
  "slides": [
    {{
      "layout": "content",
      "title": "页面标题",
      "bullets_paragraph": ["要点1段落形式（20~40字，完整的一句话描述）", "要点2段落形式", "要点3段落形式"],
      "bullets": ["要点1精炼形式（5~15字，关键词短语）", "要点2精炼形式", "要点3精炼形式"]
    }},
    {{
      "layout": "two_column",
      "title": "页面标题",
      "left_paragraph": ["左栏要点1段落形式（20~40字）", "左栏要点2段落形式"],
      "left": ["左栏要点1精炼（5~15字）", "左栏要点2精炼"],
      "right_paragraph": ["右栏要点1段落形式", "右栏要点2段落形式"],
      "right": ["右栏要点1精炼", "右栏要点2精炼"]
    }},
    {{
      "layout": "cards",
      "title": "页面标题",
      "cards": [
        {{"number": "90%", "label": "指标名称（5~15字）", "desc": "说明文字（20~50字）"}},
        {{"number": "3x", "label": "效率提升", "desc": "说明文字"}}
      ]
    }},
    {{
      "layout": "section",
      "title": "章节标题（5~15字）"
    }},
    {{
      "layout": "code",
      "title": "标题",
      "code": "代码或流程文本",
      "bullets": ["说明1（15~30字）", "说明2"]
    }},
    {{
      "layout": "quote",
      "quote": "引言或核心观点（20~40字）",
      "source": "出处",
      "bullets": ["补充要点1（15~30字）"]
    }},
    {{
      "layout": "timeline",
      "title": "流程标题",
      "steps": [
        {{"label": "步骤名（5~15字）", "desc": "描述（20~50字）"}},
        {{"label": "步骤名", "desc": "描述"}}
      ]
    }}
  ]
}}

## 内容填充规则
1. **严格遵循布局大纲**：大纲中每页是什么 layout，输出就用什么 layout，不要改变
2. **从手稿提取**：所有文字必须来自手稿原文，不要编造
3. **双编码要求**（content 和 two_column 布局）：
   - bullets_paragraph：段落形式，20~40个字，完整的一句话描述，包含数据和因果
   - bullets：要点形式，5~15个字，精炼的关键词短语
4. **元素字符限制**：不超过上方布局元素规格中的 suggested_chars
5. **严禁 emoji，所有文本中文**
6. **保留手稿中的完整信息**，不要过度精简

请直接输出 JSON:""")


# ─── HTML 预览（前端展示用）───

def _e(text: str) -> str:
    """HTML 转义"""
    import html as html_mod
    return html_mod.escape(str(text)) if text else ""


def _build_html_preview(slides_data: List[Dict], title: str, subtitle: str, theme: dict = None) -> str:
    t = theme or _DEFAULT_THEME
    acc = t.get("acc", "#7C3AED")
    pri = t.get("pri", "#1B2A4A")
    pri2 = t.get("pri2", "#2d3f66")
    txt2 = t.get("txt2", "#64748B")
    acc2 = t.get("acc2", "#A78BFA")
    acc3 = t.get("acc3", "#EDE9FE")
    bg2 = t.get("bg2", "#F8F7FF")
    bg3 = t.get("bg3", "#F1F5F9")

    # 提取 acc 的 RGB 分量用于 rgba()
    ah = acc.lstrip("#")
    ar, ag, ab = int(ah[0:2], 16), int(ah[2:4], 16), int(ah[4:6], 16)

    cards = f"""<div class="card cover"><div class="num">0</div>
      <h2>{_e(title)}</h2><p>{_e(subtitle)}</p></div>"""
    for i, s in enumerate(slides_data, 1):
        lo = s.get("layout", "content")
        st = _e(s.get("title", ""))
        if lo == "section":
            cards += f'<div class="card sect"><div class="num">{i}</div><h2>{st}</h2></div>'
        elif lo == "two_column":
            lh = "".join(f"<li>{_e(x)}</li>" for x in s.get("left", []))
            rh = "".join(f"<li>{_e(x)}</li>" for x in s.get("right", []))
            cards += f'<div class="card"><div class="num">{i}</div><h3>{st}</h3><div class="cols"><ul>{lh}</ul><ul>{rh}</ul></div></div>'
        elif lo == "cards":
            ch = ""
            for c in s.get("cards", []):
                ch += f'<div class="kpi"><b>{_e(c.get("number",""))}</b><span>{_e(c.get("label",""))}</span></div>'
            cards += f'<div class="card"><div class="num">{i}</div><h3>{st}</h3><div class="kpis">{ch}</div></div>'
        elif lo == "timeline":
            sh = ""
            for j, step in enumerate(s.get("steps", [])):
                sh += f'<div class="step"><b>{j+1}</b><span>{_e(step.get("label",""))}</span><small>{_e(step.get("desc",""))}</small></div>'
            cards += f'<div class="card"><div class="num">{i}</div><h3>{st}</h3><div class="steps">{sh}</div></div>'
        elif lo == "code":
            code_h = _e(s.get("code", "")).replace("\n", "<br>")
            bh = "".join(f"<li>{_e(x)}</li>" for x in s.get("bullets", []))
            cards += f'<div class="card"><div class="num">{i}</div><h3>{st}</h3><pre>{code_h}</pre><ul>{bh}</ul></div>'
        elif lo == "quote":
            q = _e(s.get("quote", ""))
            bh = "".join(f"<li>{_e(x)}</li>" for x in s.get("bullets", []))
            cards += f'<div class="card"><div class="num">{i}</div><blockquote>{q}</blockquote><ul>{bh}</ul></div>'
        elif lo == "toc":
            ih = "".join(f"<li>{_e(x)}</li>" for x in s.get("items", []))
            cards += f'<div class="card sect"><div class="num">{i}</div><h2>目录</h2><ul>{ih}</ul></div>'
        else:
            bh = "".join(f"<li>{_e(x)}</li>" for x in s.get("bullets", []))
            cards += f'<div class="card"><div class="num">{i}</div><h3>{st}</h3><ul>{bh}</ul></div>'
    total = len(slides_data) + 2
    return f"""<!DOCTYPE html><html lang="zh-CN"><head><meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>{_e(title)}</title>
<style>
*{{box-sizing:border-box;margin:0;padding:0}}
body{{font-family:'Microsoft YaHei',sans-serif;background:linear-gradient(135deg,{acc3} 0%,{bg2} 50%,#f0f9ff 100%);
  min-height:100vh;padding:32px 20px;color:#1E293B}}
.c{{max-width:860px;margin:0 auto}}
header{{text-align:center;margin-bottom:32px}}
header h1{{font-size:28px;font-weight:700;color:{pri};margin-bottom:8px}}
header p{{color:{txt2};font-size:15px}}
.bar{{display:flex;justify-content:center;margin-bottom:24px}}
.cnt{{font-size:14px;color:{acc};font-weight:600;background:rgba({ar},{ag},{ab},.08);padding:6px 14px;border-radius:20px}}
.card{{background:#fff;border-radius:16px;padding:28px 32px;margin-bottom:20px;
  box-shadow:0 2px 12px rgba(0,0,0,.04);border:1px solid rgba({ar},{ag},{ab},.06);position:relative}}
.num{{position:absolute;top:12px;right:16px;font-size:12px;color:#94a3b8;font-weight:600;
  background:{bg3};padding:2px 10px;border-radius:10px}}
.cover{{background:linear-gradient(135deg,{pri},{pri2});color:#fff;text-align:center;padding:48px 32px}}
.cover .num{{background:rgba(255,255,255,.15);color:rgba(255,255,255,.7)}}
.cover h2{{font-size:32px;font-weight:700;margin-bottom:12px}}
.cover p{{font-size:16px;color:rgba(255,255,255,.7)}}
.sect{{background:linear-gradient(135deg,{acc3},{bg2});text-align:center;padding:48px 32px;border:2px solid rgba({ar},{ag},{ab},.15)}}
.sect .num{{background:rgba({ar},{ag},{ab},.1);color:{acc}}}
.sect h2{{font-size:28px;font-weight:700;color:{acc}}}
h3{{font-size:20px;font-weight:600;color:{pri};margin-bottom:16px;padding-bottom:10px;border-bottom:2px solid {acc3}}}
ul{{list-style:none;padding:0}}
li{{padding:8px 0 8px 20px;position:relative;font-size:15px;line-height:1.7;color:#334155}}
li::before{{content:'\\25cf';position:absolute;left:0;color:{acc};font-size:10px;top:12px}}
.cols{{display:flex;gap:24px}}.cols ul{{flex:1}}
pre{{background:#1e1e2e;color:#cdd6f4;padding:16px;border-radius:10px;font-size:14px;line-height:1.7;
  overflow-x:auto;margin-bottom:12px;font-family:Consolas,monospace}}
blockquote{{font-size:18px;font-style:italic;color:{pri};border-left:4px solid {acc};
  padding:12px 20px;margin-bottom:16px;background:{bg2};border-radius:0 8px 8px 0}}
.kpis{{display:flex;gap:16px;flex-wrap:wrap}}
.kpi{{flex:1;min-width:120px;background:{bg2};border-radius:12px;padding:20px 16px;text-align:center}}
.kpi b{{display:block;font-size:28px;font-weight:800;color:{acc};margin-bottom:4px}}
.kpi span{{font-size:14px;color:{txt2}}}
.steps{{display:flex;gap:12px;flex-wrap:wrap}}
.step{{flex:1;min-width:100px;text-align:center;padding:16px 8px}}
.step b{{display:block;width:32px;height:32px;border-radius:50%;background:{acc};color:#fff;
  font-size:14px;line-height:32px;margin:0 auto 8px}}
.step span{{display:block;font-size:14px;font-weight:600;color:{pri};margin-bottom:4px}}
.step small{{font-size:12px;color:{txt2}}}
</style></head><body><div class="c">
<header><h1>{_e(title)}</h1><p>{_e(subtitle)}</p></header>
<div class="bar"><span class="cnt">共 {total} 页</span></div>
{cards}
</div></body></html>"""


# ─── 功能页面自动插入（参照 PPTAgent 的 _add_functional_layouts）───

def _resolve_dual_encoding(slides_data: list) -> list:
    """双编码后处理：根据布局类型选择段落形式或要点形式"""
    for sd in slides_data:
        layout = sd.get("layout", "content")

        if layout == "content":
            # 优先使用段落形式（更完整），如果存在的话
            if "bullets_paragraph" in sd:
                sd["bullets"] = sd.pop("bullets_paragraph")
        elif layout == "two_column":
            # 左右栏都优先使用段落形式
            if "left_paragraph" in sd:
                sd["left"] = sd.pop("left_paragraph")
            if "right_paragraph" in sd:
                sd["right"] = sd.pop("right_paragraph")
        # 其他布局（cards, timeline, code, quote）保持原样，不做转换

        # 清理可能残留的 _paragraph 字段
        for key in list(sd.keys()):
            if key.endswith("_paragraph"):
                sd.pop(key, None)

    return slides_data


def _ensure_functional_slides(slides_data: list, outline_slides: list, ppt_title: str) -> list:
    """确保 TOC 和 Section 分隔页存在，自动插入缺失的"""
    result = list(slides_data)

    # 1. 确保 TOC 存在（在第一页）
    has_toc = any(s.get("layout") == "toc" for s in result[:2])
    if not has_toc:
        # 从 outline 中收集所有非 section/toc 的页面标题作为目录项
        toc_items = []
        for s in outline_slides:
            lo = s.get("layout", "content")
            if lo not in ("toc", "section"):
                summary = s.get("summary", "")
                if summary:
                    toc_items.append(summary[:20])
        if toc_items:
            result.insert(0, {"layout": "toc", "items": toc_items[:6]})
            logger.info(f"[PPT] 自动插入 TOC 页，{len(toc_items[:6])} 项")

    # 2. 确保 Section 分隔页存在（在布局类型变化处）
    # 从 outline_slides 中检测 topic 变化
    sections_to_insert = []
    prev_topic = None
    for i, os in enumerate(outline_slides):
        topic = os.get("topic", "")
        summary = os.get("summary", "")
        layout = os.get("layout", "content")
        # 当 topic 变化且前一个 topic 有内容时，插入 section 页
        if topic and topic != prev_topic and prev_topic and layout != "section":
            sections_to_insert.append((i, topic))
        if layout != "section" and layout != "toc":
            prev_topic = topic

    # 在对应位置插入 section 页
    offset = 0
    for insert_idx, topic in sections_to_insert:
        actual_idx = insert_idx + offset
        # 检查该位置附近是否已有 section 页
        nearby_has_section = any(
            result[j].get("layout") == "section"
            for j in range(max(0, actual_idx - 1), min(len(result), actual_idx + 2))
        )
        if not nearby_has_section:
            result.insert(actual_idx, {"layout": "section", "title": topic})
            offset += 1
            logger.info(f"[PPT] 自动插入 Section 页: {topic}")

    # 3. 确保结尾页由 _pptx_ending 处理，不需要在这里插入

    return result


# ─── 主入口（两阶段生成 + 原生 PPTX）───

def generate_pptx(
    module_content: str,
    title: str,
    description: str,
    user_id: int = 0,
    user_profile: dict = None,
    sse_callback=None,
) -> Dict[str, Any]:
    llm = get_resource_type_generator_llm()

    def _emit_thinking(content: str):
        if sse_callback:
            try:
                sse_callback(f'data: {json.dumps({"type": "thinking", "agent": "PPT 生成智能体", "content": content}, ensure_ascii=False)}\n\n')
            except Exception:
                pass

    def _emit_progress(content: str):
        if sse_callback:
            try:
                sse_callback(f'data: {json.dumps({"type": "progress", "content": content}, ensure_ascii=False)}\n\n')
            except Exception:
                pass

    # ═══ Stage 0: 样式分析（内容感知配色）═══
    _emit_thinking("正在分析内容主题，选择配色方案...")
    _emit_progress("正在分析配色方案...")
    logger.info("[PPT] Stage 0: 样式分析...")
    theme_result = _analyze_style(title, description, module_content, user_profile, user_id)
    _apply_theme(theme_result)
    theme_name = theme_result.get("theme_name", "默认")
    _emit_thinking(f"配色方案确定: {theme_name}主题")
    logger.info(f"[PPT] 主题: {theme_name}")

    # ═══ Stage 1: 生成内容手稿（Markdown）═══
    _emit_thinking("正在生成 PPT 内容手稿，规划每一页的讲解内容...")
    _emit_progress("正在生成内容手稿...")
    logger.info("[PPT] Stage 1: 生成内容手稿...")
    ms_prompt = MANUSCRIPT_PROMPT.substitute(
        page_count="10-15",
        module_content=module_content[:16000],
        title=title,
        description=description,
    )

    manuscript = None
    for attempt in range(3):
        try:
            manuscript = llm.chat([{"role": "system", "content": ms_prompt}])
            break
        except Exception as e:
            if attempt == 2:
                raise
            logger.warning(f"[PPT] Stage 1 LLM 调用失败，重试 {attempt+1}/3: {e}")
    record_from_mimo(llm, user_id, "pptx_manuscript")

    if not manuscript or len(manuscript) < 200:
        raise Exception("Stage 1: LLM 内容手稿生成失败或过短")
    _emit_thinking(f"内容手稿生成完成，共 {len(manuscript)} 字，正在规划页面布局...")
    logger.info(f"[PPT] 手稿生成完成，{len(manuscript)} 字")

    # ═══ Stage 2a: 布局选择（参照 PPTAgent 的 layout_selector）═══
    _emit_progress("正在规划页面布局...")
    logger.info("[PPT] Stage 2a: 布局选择...")
    layouts_text = _build_layouts_text()
    outline_prompt = OUTLINE_PROMPT.substitute(
        manuscript=manuscript[:12000],
        layouts_text=layouts_text,
    )

    outline_result = None
    for attempt in range(3):
        try:
            outline_result = llm.chat_json([{"role": "system", "content": outline_prompt}])
            break
        except Exception as e:
            if attempt == 2:
                raise
            logger.warning(f"[PPT] Stage 2a LLM 调用失败，重试 {attempt+1}/3: {e}")
    record_from_mimo(llm, user_id, "pptx_outline")

    ppt_title = outline_result.get("title", title)
    ppt_subtitle = outline_result.get("subtitle", description)
    outline_slides = outline_result.get("slides", [])
    if not outline_slides:
        raise Exception("Stage 2a: LLM 未生成任何布局大纲")
    layout_names = [s.get('layout', 'content') for s in outline_slides]
    _emit_thinking(f"布局规划完成，共 {len(outline_slides)} 页，正在搜索配图...")
    logger.info(f"[PPT] 布局选择完成，{len(outline_slides)} 页: {layout_names}")

    # ═══ Stage 2b: 图片搜索（在内容生成之前，让 LLM 知道有什么图可用）═══
    _emit_progress("正在搜索配图...")
    logger.info("[PPT] Stage 2b: 图片搜索...")
    image_queries = _generate_image_queries(outline_slides, llm)
    if image_queries:
        outline_slides = _search_images_for_slides(outline_slides, image_queries, _emit_thinking)
        record_from_mimo(llm, user_id, "pptx_image_query")
    else:
        logger.info("[PPT] Stage 2b: 无需搜索配图")

    # 构建配图信息，传给 CONTENT_PROMPT
    image_info_lines = []
    for idx, s in enumerate(outline_slides):
        if "_image_data" in s:
            image_info_lines.append(f"- 第{idx+1}页 [{s.get('layout')}]: 已有配图（{_image_query_desc(s)}）")
    if image_info_lines:
        image_info = "\n\n## 可用配图\n以下页面已有配图，填充内容时文字要精简（控制在 3 条要点以内），为图片留出空间：\n" + "\n".join(image_info_lines)
    else:
        image_info = ""

    # ═══ Stage 2c: 内容填充（参照 PPTAgent 的 editor）═══
    _emit_progress("正在填充页面内容...")
    logger.info("[PPT] Stage 2c: 内容填充...")
    # 序列化前去除内部字段（_image_data 是 BytesIO，不可序列化）
    outline_for_prompt = {k: v for k, v in outline_result.items() if k != "slides"}
    outline_for_prompt["slides"] = [
        {k: v for k, v in s.items() if not k.startswith("_")}
        for s in outline_slides
    ]
    content_prompt = CONTENT_PROMPT.substitute(
        outline_json=json.dumps(outline_for_prompt, ensure_ascii=False, indent=2),
        manuscript=manuscript[:12000],
        layouts_text=layouts_text,
        image_info=image_info,
    )

    content_result = None
    for attempt in range(3):
        try:
            content_result = llm.chat_json([{"role": "system", "content": content_prompt}])
            break
        except Exception as e:
            if attempt == 2:
                raise
            logger.warning(f"[PPT] Stage 2b LLM 调用失败，重试 {attempt+1}/3: {e}")
    record_from_mimo(llm, user_id, "pptx_content")

    slides_data = content_result.get("slides", [])
    if not slides_data:
        raise Exception("Stage 2b: LLM 未生成任何 slide 内容")
    _emit_thinking(f"内容填充完成，共 {len(slides_data)} 页 slide，正在优化排版...")
    logger.info(f"[PPT] 内容填充完成，{len(slides_data)} 页 slide")

    # ═══ 将图片数据从 outline_slides 传递到 slides_data（此时索引 1:1 对应）═══
    for idx, sd in enumerate(slides_data):
        if idx < len(outline_slides) and "_image_data" in outline_slides[idx]:
            sd["_image_data"] = outline_slides[idx]["_image_data"]
            sd["_image_query"] = outline_slides[idx].get("_image_query", "")

    # ═══ 双编码后处理 ═══
    slides_data = _resolve_dual_encoding(slides_data)
    logger.info(f"[PPT] 双编码解析完成")

    # ═══ 功能页面自动插入 ═══
    slides_data = _ensure_functional_slides(slides_data, outline_slides, ppt_title)
    logger.info(f"[PPT] 功能页面处理完成，{len(slides_data)} 页 slide")

    # ═══ 验证压缩（参照 PPTAgent 的 length_rewrite）═══
    slides_data = _validate_and_compress(slides_data, llm=llm)
    logger.info(f"[PPT] 验证压缩完成")

    # ═══ 用 python-pptx 生成原生 PPTX ═══
    _emit_thinking("正在生成 PPT 文件...")
    _emit_progress("正在生成 PPT 文件...")
    logger.info("[PPT] 生成原生 PPTX...")
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    # 封面
    _pptx_title(prs, ppt_title, ppt_subtitle)

    # 内容页（带视觉变体）
    for idx, sd in enumerate(slides_data):
        layout = sd.get("layout", "content")
        builder = _PPTX_BUILDERS.get(layout)
        if builder:
            variant = idx % _VARIANT_COUNT.get(layout, 1)
            builder(prs, sd, variant)
        else:
            _pptx_content(prs, {"title": sd.get("title", ""), "bullets": sd.get("bullets", [])}, 0)

    # 结尾
    _pptx_ending(prs, ppt_title)

    # 保存
    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    filename = f"{uuid.uuid4().hex[:12]}.pptx"
    (PPTX_DIR / filename).write_bytes(pptx_bytes)
    logger.info(f"[PPT] 已保存: {filename} ({len(pptx_bytes)} bytes)")

    # 上传到七牛云
    _emit_thinking(f"PPT 文件生成完成 ({len(pptx_bytes) // 1024}KB)，正在上传...")
    from app.services.qiniu_client import qiniu_client
    pptx_url = ""
    try:
        pptx_url = qiniu_client.upload_bytes(pptx_bytes, filename, "pptx")
        _emit_thinking(f"上传完成，共 {len(slides_data) + 2} 页幻灯片")
        logger.info(f"[PPT] 已上传七牛云: {pptx_url}")
    except Exception as e:
        logger.warning(f"[PPT] 七牛云上传失败: {e}")

    return {
        "module_type": "pptx",
        "title": ppt_title,
        "description": ppt_subtitle,
        "content": _build_html_preview(slides_data, ppt_title, ppt_subtitle, theme_result),
        "pptx_filename": filename,
        "pptx_url": pptx_url,
        "slide_count": len(slides_data) + 2,
    }
