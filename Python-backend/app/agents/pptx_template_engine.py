"""
PPTX 模板引擎 — 基于 PPTAgent 模板的幻灯片生成
从预设计的模板 PPTX 中克隆幻灯片并替换文字，保留所有视觉设计元素（渐变、形状、图片等）
"""
import json
import copy
import logging
import re
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn

logger = logging.getLogger("agents.pptx_template_engine")

TEMPLATES_DIR = Path(__file__).resolve().parent.parent.parent / "templates"

# 我们的布局类型 → 模板布局类型的映射（按优先级排序）
# 使用模糊匹配：re.search(pattern, layout_name)
LAYOUT_MAPPING = {
    "content": [
        r"Text Header.*Bulleted.*:text",         # default: Title Text Header with Multi-level Bulleted...
        r"Banner.*Bullet.*:text",                 # thu: Wide Header Banner with Icon Above Bullet Point
        r"Banner Title.*Bulleted.*:text",         # beamer: Single Top Banner Title with Left-Aligned Bulleted
        r"Bulleted.*List.*:text",                 # 通用：带要点的文本布局
        r"Banner:text",
        r"Bullet Point:text",
    ],
    "two_column": [
        r"Side-by-Side.*Column.*:text",           # thu: Two Side-by-Side Text Columns...
        r"Two.*Column:text",
        r"Parallel.*Column:text",
    ],
    "cards": [
        r"Grid.*Thumbnail.*:image",               # thu: Grid of Slide Thumbnails...
        r"grid table.*:image",                    # beamer: Title bar, grid table...
        r"Grid.*Diagram.*:image",                 # default: Centered Grid Diagram...
    ],
    "timeline": [
        r"Bullet.*List:text",
        r"Bulleted.*List.*:text",
    ],
    "quote": [
        r"Full-Slide.*Text.*:text",               # default: Single Full-Slide Text Block...
        r"Text Block.*Decorative.*:text",
        r"Decorative.*Background:text",
    ],
    "section": [
        r"section outline",
    ],
    "code": [
        r"Bulleted.*List.*:text",
        r"Text Header.*:text",
    ],
    "toc": [
        r"table of contents",
    ],
    "ending": [
        r"ending",
    ],
}


class TemplateEngine:
    """基于 PPTAgent 模板的 PPT 生成引擎"""

    def __init__(self, template_name: str = "default"):
        self.template_name = template_name
        self.template_dir = TEMPLATES_DIR / template_name
        if not self.template_dir.exists():
            raise FileNotFoundError(f"模板不存在: {self.template_dir}")

        # 加载模板
        self.induction = self._load_induction()
        self.prs = Presentation(str(self.template_dir / "source.pptx"))
        self.layouts = self._parse_layouts()
        self.functional_keys = self.induction.get("functional_keys", [])

        logger.info(f"[Template] 加载模板 '{template_name}': {len(self.layouts)} 个布局, "
                    f"{len(self.prs.slides)} 张幻灯片")

    def _load_induction(self) -> dict:
        path = self.template_dir / "slide_induction.json"
        with open(path, encoding="utf-8") as f:
            return json.load(f)

    def _parse_layouts(self) -> dict:
        """解析模板布局信息"""
        layouts = {}
        for key, val in self.induction.items():
            if key in ("functional_keys", "language"):
                continue
            layouts[key] = {
                "template_id": val["template_id"],
                "slides": val.get("slides", []),
                "elements": val.get("elements", []),
            }
        return layouts

    def get_available_layouts(self) -> List[str]:
        """返回模板中可用的布局名称"""
        return list(self.layouts.keys())

    def match_layout(self, our_layout: str) -> Optional[str]:
        """将我们的布局类型映射到模板布局（正则模糊匹配）。
        如果模板没有匹配的布局，返回 None（调用方会降级为从零构建）。
        """
        candidates = LAYOUT_MAPPING.get(our_layout, [])
        for pattern in candidates:
            for key in self.layouts:
                if re.search(pattern, key, re.IGNORECASE):
                    return key
        return None

    def clone_slide(self, template_id: int) -> object:
        """从模板中深拷贝一张幻灯片（保留所有视觉元素）"""
        # template_id 是 1-based
        src_slide = self.prs.slides[template_id - 1]
        return copy.deepcopy(src_slide)

    def find_text_shapes(self, slide) -> List[Tuple[object, str]]:
        """找到幻灯片中所有包含文字的形状，返回 (shape, text) 列表"""
        text_shapes = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                full_text = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
                if full_text:
                    text_shapes.append((shape, full_text))
        return text_shapes

    def replace_text_in_slide(self, slide, elements: List[dict], content: dict):
        """替换幻灯片中的文字内容
        elements: 模板的元素定义 [{"name": "main title", "type": "text", "data": [...]}]
        content: 我们的内容 {"main title": "新标题", "main bullets": ["要点1", "要点2"]}
        """
        text_shapes = self.find_text_shapes(slide)

        for elem in elements:
            elem_name = elem["name"]
            elem_type = elem["type"]
            if elem_type != "text":
                continue
            if elem_name not in content:
                continue

            new_text = content[elem_name]
            if isinstance(new_text, list):
                new_text = "\n".join(str(t) for t in new_text)

            # 找到最匹配的形状（通过原始文字匹配）
            original_texts = elem.get("data", [])
            best_shape = self._find_best_matching_shape(text_shapes, original_texts)

            if best_shape:
                self._set_shape_text(best_shape, new_text)

    def _find_best_matching_shape(self, text_shapes, original_texts):
        """通过原始文字找到最匹配的形状"""
        if not original_texts or not text_shapes:
            return None

        best_shape = None
        best_score = 0

        for shape, shape_text in text_shapes:
            for orig in original_texts:
                # 计算相似度：共同字符比例
                orig_clean = orig.strip()[:50]
                shape_clean = shape_text.strip()[:50]
                if not orig_clean or not shape_clean:
                    continue
                # 简单匹配：检查是否有共同的关键词
                common = sum(1 for c in orig_clean if c in shape_clean)
                score = common / max(len(orig_clean), len(shape_clean))
                if score > best_score:
                    best_score = score
                    best_shape = shape

        # 阈值：至少 30% 相似度
        return best_shape if best_score > 0.3 else None

    def _set_shape_text(self, shape, new_text: str):
        """设置形状的文字内容，保留原始格式"""
        tf = shape.text_frame
        paragraphs = new_text.split("\n")

        # 保留第一个段落的格式作为模板
        if not tf.paragraphs:
            return

        # 获取原始格式
        first_para = tf.paragraphs[0]
        orig_font_size = first_para.font.size
        orig_font_bold = first_para.font.bold
        orig_font_name = first_para.font.name
        orig_font_color = None
        try:
            orig_font_color = first_para.font.color.rgb
        except Exception:
            pass

        # 清除所有现有段落
        for p in tf.paragraphs:
            p.clear()

        # 删除多余的段落（保留第一个）
        while len(tf.paragraphs) > 1:
            p_elem = tf.paragraphs[-1]._p
            p_elem.getparent().remove(p_elem)

        # 写入新文字
        for i, line in enumerate(paragraphs):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # 清除旧的 runs
            for run in p.runs:
                p._p.remove(run._r)

            # 创建新 run
            run = p.add_run()
            run.text = line.strip()

            # 应用原始格式
            if orig_font_size:
                run.font.size = orig_font_size
            if orig_font_bold:
                run.font.bold = orig_font_bold
            if orig_font_name:
                run.font.name = orig_font_name
            if orig_font_color:
                try:
                    run.font.color.rgb = orig_font_color
                except Exception:
                    pass

    def generate_slides(self, slides_data: List[dict], target_prs: Presentation):
        """为每页数据从模板生成幻灯片，添加到 target_prs 中"""
        import copy as copy_mod

        for idx, sd in enumerate(slides_data):
            layout = sd.get("layout", "content")

            # 映射到模板布局
            template_layout = self.match_layout(layout)
            if not template_layout or template_layout not in self.layouts:
                logger.warning(f"[Template] 无法映射布局 '{layout}'，跳过")
                continue

            layout_info = self.layouts[template_layout]
            template_id = layout_info["template_id"]
            elements = layout_info["elements"]

            # 准备内容映射
            content = self._prepare_content(layout, sd)

            # 深拷贝模板幻灯片
            src_slide = self.prs.slides[template_id - 1]
            cloned_slide_xml = copy.deepcopy(src_slide._element)

            # 添加到目标 presentation（通过操作 XML）
            self._add_slide_to_prs(cloned_slide_xml, src_slide, target_prs)

            # 获取刚添加的幻灯片
            new_slide = target_prs.slides[-1]

            # 替换文字
            self.replace_text_in_slide(new_slide, elements, content)

            logger.debug(f"[Template] 页 {idx}: {layout} -> {template_layout}")

    def _prepare_content(self, layout: str, sd: dict) -> dict:
        """将我们的 slide 数据转换为模板元素名称→内容的映射"""
        content = {}

        if layout == "content":
            content["main title"] = sd.get("title", "")
            bullets = sd.get("bullets", [])
            content["quote and explanation"] = "\n".join(str(b) for b in bullets[:5])
            content["main bullets"] = "\n".join(f"• {b}" for b in bullets[:5])
            content["main text block"] = "\n".join(str(b) for b in bullets[:5])
            content["main paragraph"] = "\n".join(str(b) for b in bullets[:5])

        elif layout == "two_column":
            content["main title"] = sd.get("title", "")
            left = sd.get("left", [])
            right = sd.get("right", [])
            content["left section title"] = "对比 A"
            content["left section bullets"] = "\n".join(f"• {b}" for b in left[:4])
            content["right section title"] = "对比 B"
            content["right section bullets"] = "\n".join(f"• {b}" for b in right[:4])

        elif layout == "cards":
            content["main title"] = sd.get("title", "")
            cards = sd.get("cards", [])
            lines = []
            for c in cards[:4]:
                num = c.get("number", "")
                label = c.get("label", "")
                desc = c.get("desc", "")
                lines.append(f"【{num}】{label}：{desc}")
            content["quote and explanation"] = "\n".join(lines)
            content["main bullets"] = "\n".join(f"• {l}" for l in lines)
            content["main body text"] = "\n".join(lines)

        elif layout == "timeline":
            content["main title"] = sd.get("title", "")
            steps = sd.get("steps", [])
            lines = []
            for i, s in enumerate(steps[:5]):
                label = s.get("label", "")
                desc = s.get("desc", "")
                lines.append(f"{i+1}. {label}：{desc}")
            content["quote and explanation"] = "\n".join(lines)
            content["main bullets"] = "\n".join(f"• {l}" for l in lines)
            content["main text block"] = "\n".join(lines)

        elif layout == "quote":
            quote = sd.get("quote", "")
            source = sd.get("source", "")
            bullets = sd.get("bullets", [])
            text = f"“{quote}”"
            if source:
                text += f"\n—— {source}"
            if bullets:
                text += "\n" + "\n".join(f"• {b}" for b in bullets[:3])
            content["main paragraph"] = text
            content["quote and explanation"] = text

        elif layout == "section":
            content["main title"] = sd.get("title", "")
            content["main body"] = sd.get("title", "")
            content["main paragraph"] = sd.get("title", "")

        elif layout == "toc":
            items = sd.get("items", [])
            content["main title"] = "目录"
            content["bulleted list"] = "\n".join(f"• {item}" for item in items[:6])
            content["section list"] = "\n".join(f"• {item}" for item in items[:6])
            content["main bullets"] = "\n".join(f"• {item}" for item in items[:6])

        elif layout == "code":
            content["main title"] = sd.get("title", "")
            code = sd.get("code", "")
            bullets = sd.get("bullets", [])
            text = code[:500]
            if bullets:
                text += "\n" + "\n".join(f"• {b}" for b in bullets[:3])
            content["quote and explanation"] = text
            content["main bullets"] = "\n".join(f"• {b}" for b in bullets[:4])

        return content

    def _add_slide_to_prs(self, slide_xml, src_slide, target_prs):
        """将克隆的幻灯片 XML 添加到目标 presentation（通过空白页替换内容）"""
        # 使用空白布局添加一张新幻灯片
        slide_layout = target_prs.slide_layouts[6]  # blank layout
        new_slide = target_prs.slides.add_slide(slide_layout)

        # 用克隆的 XML 替换空白幻灯片的 spTree（形状树）
        new_cSld = new_slide._element.find(qn('p:cSld'))
        cloned_cSld = slide_xml.find(qn('p:cSld'))

        if new_cSld is not None and cloned_cSld is not None:
            new_sp_tree = new_cSld.find(qn('p:spTree'))
            cloned_sp_tree = cloned_cSld.find(qn('p:spTree'))

            if new_sp_tree is not None and cloned_sp_tree is not None:
                # 清除空白页的所有形状
                for child in list(new_sp_tree):
                    new_sp_tree.remove(child)
                # 复制模板页的所有形状
                for child in cloned_sp_tree:
                    new_sp_tree.append(copy.deepcopy(child))

            # 复制背景（如果模板有自定义背景）
            cloned_bg = cloned_cSld.find(qn('p:bg'))
            if cloned_bg is not None:
                existing_bg = new_cSld.find(qn('p:bg'))
                if existing_bg is not None:
                    new_cSld.remove(existing_bg)
                new_cSld.insert(0, copy.deepcopy(cloned_bg))


def get_template_engine(template_name: str = "default") -> TemplateEngine:
    """获取模板引擎实例"""
    return TemplateEngine(template_name)


def list_templates() -> List[dict]:
    """列出所有可用模板"""
    templates = []
    if TEMPLATES_DIR.exists():
        for d in TEMPLATES_DIR.iterdir():
            if d.is_dir() and (d / "source.pptx").exists():
                desc_file = d / "description.txt"
                desc = desc_file.read_text(encoding="utf-8").strip() if desc_file.exists() else ""
                templates.append({"name": d.name, "description": desc})
    return templates
